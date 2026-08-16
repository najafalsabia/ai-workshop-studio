import os
import sys
import json
import copy
import tempfile
from concurrent.futures import ThreadPoolExecutor, as_completed
import zipfile
import io
from pathlib import Path

# Reconfigure stdout/stderr to support printing Arabic Unicode characters and emojis in Windows shell
sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')

from fastapi import FastAPI, HTTPException, Body, UploadFile, File
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles

# Setup python path to load Backend modules
root_dir = Path(__file__).resolve().parent
backend_dir = root_dir / "Backend"
sys.path.insert(0, str(backend_dir))

from dotenv import load_dotenv
load_dotenv(backend_dir / ".env")

app = FastAPI(title="Capsule Studio API Server")

# Import Backend Functions
try:
    from idea_agent import generate_titles
    from plan_builder import build_plan
    from content_generator import generate_content_with_sources
    from activities_generator import extract_lab_contexts, suggest_lab_types, generate_labs
    from quiz_generator import generate_quiz
    from pptx_export import export_to_pptx
    from quality_checklist import run_quality_checklist
    from notebook_builder import build_notebook, save_all_lab_notebooks
    from quiz_doc_builder import save_quiz_docx
    from kahoot_export import export_quiz_to_kahoot_xlsx
    from llm_client import ask_llm_for_json
    from feedback_loop import build_revision_prompt
    import nbformat
except ImportError as e:
    print(f"Error importing Backend modules: {e}")
    sys.exit(1)


TRANSLATE_BATCH_SIZE = 5  # slides per translation call — keeps each reply well under the output token ceiling


def _translate_json_blob(data, context: str):
    """One translation call for a JSON-serializable payload small enough
    to safely fit in one reply. Raises on failure — callers decide how to
    handle that (see api_translate below)."""
    prompt = f"""You are a professional technical translator translating workshop materials into Arabic.
Translate the following JSON content from English to Arabic.
Keep the JSON keys, structure, types, and values exactly the same. Only translate the text values.
For code blocks or code cells in notebooks, do NOT translate the actual code itself, only translate comments and explanations.
Ensure technical accuracy and natural phrasing.

Context: {context}

JSON Content:
{json.dumps(data, indent=2, ensure_ascii=False)}

Reply with ONLY the translated JSON, nothing else."""
    return ask_llm_for_json(prompt)


@app.post("/api/translate")
async def api_translate(payload: dict = Body(...)):
    data = payload.get("data")
    context = payload.get("context", "Workshop material")
    if not data:
        raise HTTPException(status_code=400, detail="Missing data to translate")

    def _translate_list_field(list_key: str, batch_size: int):
        """
        Shared batching logic for any top-level list field (slides, labs)
        that might be too big to translate in one call. Runs batches in
        parallel (bounded), reassembles in ORIGINAL order regardless of
        completion order, and falls back to the original English for any
        batch that fails or comes back the wrong size — a bad/failed batch
        never takes down the rest of the translation.
        """
        items = data[list_key]
        batch_starts = list(range(0, len(items), batch_size))

        def _translate_one_batch(batch_start: int):
            batch_end = min(batch_start + batch_size, len(items))
            batch = items[batch_start:batch_end]
            try:
                batch_result = _translate_json_blob({list_key: batch}, context)
                batch_items = batch_result.get(list_key, [])
            except Exception as e:
                print(f"⚠️  Translation batch {batch_start}-{batch_end} of '{list_key}' failed ({e}) — using original English for this batch.")
                batch_items = batch

            if len(batch_items) != len(batch):
                print(f"⚠️  Translation batch {batch_start}-{batch_end} of '{list_key}' returned {len(batch_items)} items, expected {len(batch)} — using original English for this batch.")
                batch_items = batch
            return batch_start, batch_items

        max_workers = min(3, len(batch_starts)) if batch_starts else 1
        results_by_start = {}
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = [executor.submit(_translate_one_batch, start) for start in batch_starts]
            for future in as_completed(futures):
                start, batch_items = future.result()
                results_by_start[start] = batch_items

        translated_items = []
        for start in batch_starts:  # original order, not completion order
            translated_items.extend(results_by_start[start])
        return translated_items

    try:
        # A full slide deck (title + speaker_notes + sources per slide,
        # ~20-40 slides on a longer workshop) — or a set of labs, where
        # EACH coding lab carries two full notebooks worth of markdown +
        # code — can be big enough that ONE translation call's reply gets
        # cut off before the JSON finishes. That truncated, invalid JSON
        # is what was causing this endpoint's intermittent 500s (first
        # found on slides, then again on labs — same root cause, just a
        # different field). Fix: translate both in small batches instead
        # of the whole thing in one shot, then reassemble. Every other
        # payload shape (plan, quiz — much smaller) still goes through as
        # a single call, unchanged.
        if isinstance(data, dict) and isinstance(data.get("slides"), list) and len(data["slides"]) > TRANSLATE_BATCH_SIZE:
            translated = {**data, "slides": _translate_list_field("slides", TRANSLATE_BATCH_SIZE)}
        elif isinstance(data, dict) and isinstance(data.get("labs"), list) and len(data["labs"]) > 1:
            # One lab per batch — a single coding lab (two full notebooks
            # of cells) is already large enough to risk truncation on its
            # own, so labs get a tighter batch size than slides.
            translated = {**data, "labs": _translate_list_field("labs", 1)}
        else:
            translated = _translate_json_blob(data, context)

        return JSONResponse(content=translated)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate_titles")
async def api_generate_titles(payload: dict = Body(...)):
    audience = payload.get("audience", "")
    age = payload.get("age", "")
    duration = payload.get("duration", "")
    goal = payload.get("goal", "")
    notes = payload.get("notes", "")
    uploaded_content = payload.get("uploaded_content", "")

    # "I already have an idea" mode: the user's own idea IS the title.
    # Calling the LLM here would silently ignore what they typed and hand
    # back unrelated AI-suggested titles instead — so we skip generation
    # entirely and echo their idea back in the same response shape the
    # frontend already expects, no frontend changes needed.
    idea_mode = payload.get("idea_mode", "")
    idea_input = payload.get("idea_input", "").strip()
    if idea_mode == "have" and idea_input:
        return JSONResponse(content={
            "titles": [{"title": idea_input, "why": "This is the idea you provided."}]
        })

    if uploaded_content:
        notes = f"Uploaded Source Material:\n{uploaded_content}\n\n" + notes
    
    # Inject Field/Domain if supplied
    field = payload.get("field", "")
    if field:
        notes = f"Workshop Field/Domain: {field}\n" + notes
        
    try:
        titles = generate_titles(
            audience=audience,
            age=age,
            duration=duration,
            goal=goal,
            notes=notes
        )
        return JSONResponse(content=titles)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/build_plan")
async def api_build_plan(payload: dict = Body(...)):
    title = payload.get("title", "")
    audience = payload.get("audience", "")
    age = payload.get("age", "")
    duration = payload.get("duration", "")
    goal = payload.get("goal", "")
    notes = payload.get("notes", "")
    uploaded_content = payload.get("uploaded_content", "")
    use_cycle = payload.get("use_cycle", False)
    
    if uploaded_content:
        notes = f"Uploaded Source Material:\n{uploaded_content}\n\n" + notes

    try:
        if use_cycle:
            # Real, structured parameters — this is what actually triggers
            # plan_builder's deterministic explain/lab/break skeleton math
            # and role-tagging, not a text hint the model might ignore.
            plan = build_plan(
                title=title,
                audience=audience,
                age=age,
                duration=duration,
                goal=goal,
                notes=notes,
                use_lab_cycle=True,
                explain_minutes=int(payload.get("explain_min", 30)),
                lab_minutes=int(payload.get("lab_min", 20)),
                break_minutes=int(payload.get("break_min", 15)),
            )
        else:
            plan = build_plan(
                title=title,
                audience=audience,
                age=age,
                duration=duration,
                goal=goal,
                notes=notes,
            )
        return JSONResponse(content=plan)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/parse_file")
async def api_parse_file(file: UploadFile = File(...)):
    filename = file.filename.lower()
    contents = await file.read()
    
    if filename.endswith(".txt") or filename.endswith(".json"):
        return JSONResponse(content={"text": contents.decode("utf-8", errors="ignore")})
    elif filename.endswith(".docx"):
        import io
        try:
            from docx import Document
            doc = Document(io.BytesIO(contents))
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return JSONResponse(content={"text": "\n".join(full_text)})
        except Exception as docx_err:
            raise HTTPException(status_code=500, detail=f"Failed to parse DOCX: {docx_err}")
    elif filename.endswith(".pptx"):
        import io
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(contents))
            slide_texts = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                texts_in_slide = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                texts_in_slide.append(line)
                if texts_in_slide:
                    slide_texts.append(f"--- Slide {slide_num} ---\n" + "\n".join(texts_in_slide))
            full_text = "\n\n".join(slide_texts)
            return JSONResponse(content={"text": full_text})
        except ImportError:
            raise HTTPException(status_code=500, detail="python-pptx is not installed. Run: pip install python-pptx")
        except Exception as pptx_err:
            raise HTTPException(status_code=500, detail=f"Failed to parse PPTX: {pptx_err}")
    else:
        raise HTTPException(status_code=400, detail="Unsupported file format. Please upload .txt, .json, .docx, or .pptx")


@app.post("/api/generate_content")
async def api_generate_content(payload: dict = Body(...)):
    title = payload.get("title", "")
    learning_objectives = payload.get("learning_objectives", [])
    outline = payload.get("outline", [])
    style = payload.get("style", "Clean & Minimal")
    section_styles = payload.get("section_styles", {})
    mimic_example = payload.get("mimic_example", "")
    uploaded_content = payload.get("uploaded_content", "")

    try:
        content = generate_content_with_sources(
            title=title,
            learning_objectives=learning_objectives,
            outline=outline,
            mimic_example=mimic_example,
            default_style=style,
            section_styles=section_styles,
            uploaded_content=uploaded_content
        )
        return JSONResponse(content=content)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate_labs")
async def api_generate_labs(payload: dict = Body(...)):
    title = payload.get("title", "")
    outline = payload.get("outline") or []
    uploaded_content = payload.get("uploaded_content", "")
    content = payload.get("content")  # Step 2's slide content, if the caller sends it

    try:
        if not outline and uploaded_content:
            # Content-only mode: synthesize a single-section outline covering
            # the whole uploaded content, so extract_lab_contexts (which
            # always needs a real outline) has something to work with,
            # instead of a separate hand-rolled path with its own schema.
            outline = [{
                "section": "Uploaded Content",
                "duration_minutes": 60,
                "description": "Content provided directly by the trainer.",
            }]
            if not content:
                content = {"Uploaded Content": uploaded_content[:4000]}

        contexts = extract_lab_contexts(outline)
        if not contexts:
            return JSONResponse(content={"labs": []})

        suggestions = suggest_lab_types(title, contexts)
        # This page has no separate trainer-confirmation step for lab type
        # yet — auto-accept the AI's own suggestion as the confirmed type,
        # the same "AI suggests a sensible default" pattern used elsewhere
        # (e.g. plan_builder's role-based cycle). generate_labs still
        # requires an explicit type per lab and raises rather than
        # guessing if one is ever missing.
        confirmed_lab_types = {
            s["outline_index"]: s["lab_type"] for s in suggestions.get("suggestions", [])
        }

        result = generate_labs(
            title=title,
            lab_contexts=contexts,
            confirmed_lab_types=confirmed_lab_types,
            content=content,
        )
        return JSONResponse(content=result)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate_quiz")
async def api_generate_quiz(payload: dict = Body(...)):
    title = payload.get("title", "")
    outline = payload.get("outline") or []
    uploaded_content = payload.get("uploaded_content", "")
    content = payload.get("content")  # Step 2's slide content, if the caller sends it

    # If no outline but we have uploaded content, use it as the quiz's
    # grounding material directly (generate_quiz's content-only path).
    extra_context = ""
    if not outline and uploaded_content:
        extra_context = uploaded_content

    try:
        # question_count, if the trainer set one in the UI, overrides the
        # duration-based auto-scaling entirely (see generate_quiz's
        # question_count param in quiz_generator.py). Leave it out/None
        # to keep the automatic "scaled to teaching time" behavior.
        question_count = payload.get("question_count")
        if question_count not in (None, "", 0):
            question_count = int(question_count)
        else:
            question_count = None

        quiz = generate_quiz(
            title=title,
            outline=outline,
            content=content,
            extra_context=extra_context,
            question_count=question_count,
        )
        return JSONResponse(content=quiz)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/run_checklist")
async def api_run_checklist(payload: dict = Body(...)):
    title = payload.get("title", "AI Workshop")
    plan = payload.get("plan", {})
    content = payload.get("content", {})
    labs = payload.get("labs")
    quiz = payload.get("quiz")
    
    # Wrap labs and quiz to conform to Backend expectations if needed
    if labs and isinstance(labs, list):
        labs = {"labs": labs}
    if quiz and isinstance(quiz, list):
        quiz = {"quiz": {"questions": quiz}}
    elif quiz and isinstance(quiz, dict) and "questions" in quiz and "quiz" not in quiz:
        quiz = {"quiz": quiz}

    try:
        checklist = run_quality_checklist(
            title=title,
            plan=plan,
            content=content,
            labs_result=labs,
            quiz_result=quiz
        )
        return JSONResponse(content=checklist)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/edit_presentation")
async def api_edit_presentation(payload: dict = Body(...)):
    instruction = payload.get("instruction", "")
    deck = payload.get("deck", {})
    if not instruction or not deck:
        raise HTTPException(status_code=400, detail="Missing instruction or deck")
        
    prompt = f"""You are an expert presentation editor. The user wants to modify their presentation slides.

INSTRUCTION:
"{instruction}"

CURRENT PRESENTATION SLIDES (JSON):
{json.dumps(deck, ensure_ascii=False, indent=2)}

You can rewrite slide contents, add new slides, delete slides, or change their order.
Return the complete updated presentation as JSON with the exact same format:
{{
  "slides": [
    {{
      "slide_number": 1,
      "slide_title": "...",
      "content_type": "content_slide" or "columns_3_slide" or "timeline_slide",
      "section": "...",
      "blocks": [
         // list of blocks: heading, paragraph, bullet_list, activity, image_placeholder
      ],
      "sources": [
         // optional list of citation source dicts: author, year, title, url, exact_quote
      ],
      "slide_style": "Clean & Minimal" or "Bold & Impactful" or "Visual & Diagram-heavy" or "Data & Research" or "Interactive & Workshop"
    }},
    ...
  ]
}}

Ensure every slide you output has valid block structures following the schema:
- content_slide / roadmap_slide: block types: heading, paragraph, bullet_list, activity, image_placeholder
- columns_3_slide: block types: paragraph, columns_3 (containing a "columns" list with "heading", "text", "image_placeholder")
- timeline_slide: block types: timeline (containing "events" list with "title", "date", "text"), image_placeholder

Do NOT include any extra keys or text other than the JSON."""
    try:
        edited = ask_llm_for_json(prompt)
        # Re-sequence
        if edited and "slides" in edited:
            for idx, s in enumerate(edited["slides"], 1):
                s["slide_number"] = idx
                if "sources" not in s:
                    s["sources"] = []
        return JSONResponse(content=edited)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/refine")
async def api_refine(payload: dict = Body(...)):
    category = payload.get("category", "")
    current_data = payload.get("current_data", {})
    feedback = payload.get("feedback", "")
    
    if not category or not current_data or not feedback:
        raise HTTPException(status_code=400, detail="Missing category, current_data, or feedback")
        
    contexts = {
        "titles": "A list of workshop title suggestions, containing 'titles' which is an array of objects, each with 'title' and 'why' fields.",
        "plan": "A workshop plan, containing 'learning_objectives' (list of strings) and 'outline' (list of sections, each with section name, duration_minutes, description, and role).",
        "content": "A workshop slides content structure, containing 'slides' list where each slide has slide_number, slide_title, content_type, section, blocks, sources, slide_style.",
        "single_slide": "A single workshop slide containing slide_number, slide_title, content_type, section, blocks (a list of layout blocks like bullet_list, paragraphs, code, image_placeholder, roadmap, columns_3), sources (optional list), slide_style, speaker_notes.",
        "labs": "A list of hands-on labs containing 'labs' array where each lab has title, lab_type, duration_minutes, instructions, trainee_notebook_cells, suggested_platforms, solution_notebook_cells, instructor_notes.",
        "quiz": "A technical workshop quiz containing 'quiz' object with 'title' and 'questions' list, where each question has question, options, correct_answer, difficulty."
    }
    
    context_desc = contexts.get(category, "Workshop data structure")
    prompt = build_revision_prompt(current_data, feedback, context_desc)
    
    try:
        revised = ask_llm_for_json(prompt)
        return JSONResponse(content=revised)
    except Exception as e:
        print(f"[/api/refine ERROR] category={category!r}, feedback={feedback!r}")
        print(f"[/api/refine ERROR] {type(e).__name__}: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/export_zip")
async def api_export_zip(payload: dict = Body(...)):
    title = payload.get("title", "AI_Workshop")
    plan = payload.get("plan", {})
    content = payload.get("content", {})
    labs = payload.get("labs")
    quiz = payload.get("quiz")
    include_images = payload.get("include_images", True)
    
    try:
        # Create output files temporarily
        temp_dir = tempfile.gettempdir()
        
        # 1) Build PPTX Slides (if content has slides)
        pptx_filename = f"{title.replace(' ', '_')}.pptx"
        pptx_temp_path = os.path.join(temp_dir, pptx_filename)
        has_pptx = False
        if content and content.get("slides"):
            try:
                export_to_pptx(
                    title=title,
                    slides=content.get("slides", []),
                    output_path=pptx_temp_path,
                    include_images=include_images
                )
                has_pptx = True
            except Exception as pptx_err:
                print(f"PPTX skipped: {pptx_err}")
        
        # 2) Build Quiz DOCX (if quiz has questions)
        docx_temp_path = os.path.join(temp_dir, "ai_workshop_quiz.docx")
        has_docx = False
        quiz_data = quiz.get("quiz", quiz) if quiz else None
        if quiz_data and quiz_data.get("questions"):
            try:
                save_quiz_docx(quiz, docx_temp_path)
                has_docx = True
            except Exception as docx_err:
                print(f"DOCX skipped: {docx_err}")
        
        # 3) Build Kahoot XLSX (if docx exists)
        kahoot_temp_path = os.path.join(temp_dir, "kahoot_import.xlsx")
        kahoot_template = os.path.join(backend_dir, "KahootQuizTemplate.xlsx")
        has_kahoot = False
        if has_docx and os.path.exists(kahoot_template):
            try:
                export_quiz_to_kahoot_xlsx(quiz, kahoot_template, kahoot_temp_path)
                has_kahoot = True
            except Exception as kahoot_err:
                print(f"Kahoot skipped: {kahoot_err}")
                
        # 4) Lab Notebooks (if labs exists)
        lab_files = []
        if labs:
            labs_list = labs.get("labs", []) if isinstance(labs, dict) else labs
            for idx, lab_data in enumerate(labs_list, 1):
                try:
                    lab_type = lab_data.get("lab_type", "coding")
                    clean_t = "".join(c for c in lab_data.get("lab_title", lab_data.get("title", f"Lab_{idx}")) if c.isalnum() or c in (" ", "_", "-")).strip().replace(" ", "_")
                    if lab_type == "coding":
                        # Trainee notebook
                        nb_trainee = build_notebook(lab_data.get("trainee_notebook_cells", []))
                        nb_t_path = os.path.join(temp_dir, f"Lab_{idx}_{clean_t}_Trainee.ipynb")
                        with open(nb_t_path, "w", encoding="utf-8") as f:
                            nbformat.write(nb_trainee, f)
                        lab_files.append((nb_t_path, f"Labs/Lab_{idx}_{clean_t}_Trainee.ipynb"))
                        
                        # Instructor notebook (solution)
                        nb_sol = build_notebook(lab_data.get("solution_notebook_cells", []))
                        nb_s_path = os.path.join(temp_dir, f"Lab_{idx}_{clean_t}_Solution.ipynb")
                        with open(nb_s_path, "w", encoding="utf-8") as f:
                            nbformat.write(nb_sol, f)
                        lab_files.append((nb_s_path, f"Labs/Lab_{idx}_{clean_t}_Solution.ipynb"))
                    else:
                        # Conceptual markdown Q&A
                        md_content = f"# Lab {idx}: {lab_data.get('lab_title', lab_data.get('title', 'Conceptual Lab'))}\n\n"
                        md_content += f"## Objectives\n{lab_data.get('objectives', '')}\n\n"
                        md_content += f"## Task Description\n{lab_data.get('task_description', '')}\n\n"
                        md_content += f"## Questions\n"
                        for q in lab_data.get("questions", []):
                            md_content += f"- **Question:** {q.get('question', '')}\n"
                            md_content += f"  - **Answer Key:** {q.get('answer', '')}\n\n"
                        
                        md_path = os.path.join(temp_dir, f"Lab_{idx}_{clean_t}.md")
                        with open(md_path, "w", encoding="utf-8") as f:
                            f.write(md_content)
                        lab_files.append((md_path, f"Labs/Lab_{idx}_{clean_t}.md"))
                except Exception as lab_err:
                    print(f"Lab {idx} skipped: {lab_err}")

        # Create Zip file in memory
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            # Full JSON
            full_json = {
                "title": title,
                "plan": plan,
                "content": content,
                "labs": labs,
                "quiz": quiz
            }
            zf.writestr("ai_workshop_studio_package.json", json.dumps(full_json, indent=2, ensure_ascii=False))
            
            # Slides
            if has_pptx and os.path.exists(pptx_temp_path):
                zf.write(pptx_temp_path, f"Slides/{pptx_filename}")
                
            # Quiz
            if has_docx and os.path.exists(docx_temp_path):
                zf.write(docx_temp_path, "Quiz/ai_workshop_quiz.docx")
            if has_kahoot and os.path.exists(kahoot_temp_path):
                zf.write(kahoot_temp_path, "Quiz/kahoot_import.xlsx")
                
            # Labs
            for file_path, zip_path in lab_files:
                if os.path.exists(file_path):
                    zf.write(file_path, zip_path)

        # Cleanup temp files
        for p in [pptx_temp_path, docx_temp_path, kahoot_temp_path] + [f[0] for f in lab_files]:
            try:
                if os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass
                
        # Return Zip
        zip_buffer.seek(0)
        zip_path = os.path.join(temp_dir, f"{title.replace(' ', '_')}_package.zip")
        with open(zip_path, "wb") as f:
            f.write(zip_buffer.getvalue())
            
        return FileResponse(
            path=zip_path,
            filename=f"{title.replace(' ', '_')}_package.zip",
            media_type="application/zip"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/export_pptx")
async def api_export_pptx(payload: dict = Body(...)):
    title = payload.get("title", "AI_Workshop")
    content = payload.get("content", {})
    include_images = payload.get("include_images", True)
    
    try:
        temp_dir = tempfile.gettempdir()
        pptx_filename = f"{title.replace(' ', '_')}.pptx"
        pptx_temp_path = os.path.join(temp_dir, pptx_filename)
        
        if content and content.get("slides"):
            export_to_pptx(
                title=title,
                slides=content.get("slides", []),
                output_path=pptx_temp_path,
                include_images=include_images
            )
            return FileResponse(
                path=pptx_temp_path,
                filename=pptx_filename,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        else:
            raise HTTPException(status_code=400, detail="No slides found in content")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/export_labs")
async def api_export_labs(payload: dict = Body(...)):
    title = payload.get("title", "AI_Workshop")
    labs = payload.get("labs")
    
    try:
        temp_dir = tempfile.gettempdir()
        lab_files = []
        if labs:
            labs_list = labs.get("labs", []) if isinstance(labs, dict) else labs
            for idx, lab_data in enumerate(labs_list, 1):
                try:
                    lab_type = lab_data.get("lab_type", "coding")
                    clean_t = "".join(c for c in lab_data.get("lab_title", lab_data.get("title", f"Lab_{idx}")) if c.isalnum() or c in (" ", "_", "-")).strip().replace(" ", "_")
                    if lab_type == "coding":
                        # Trainee notebook
                        nb_trainee = build_notebook(lab_data.get("trainee_notebook_cells", []))
                        nb_t_path = os.path.join(temp_dir, f"Lab_{idx}_{clean_t}_Trainee.ipynb")
                        with open(nb_t_path, "w", encoding="utf-8") as f:
                            nbformat.write(nb_trainee, f)
                        lab_files.append((nb_t_path, f"Lab_{idx}_{clean_t}_Trainee.ipynb"))
                        
                        # Instructor notebook (solution)
                        nb_sol = build_notebook(lab_data.get("solution_notebook_cells", []))
                        nb_s_path = os.path.join(temp_dir, f"Lab_{idx}_{clean_t}_Solution.ipynb")
                        with open(nb_s_path, "w", encoding="utf-8") as f:
                            nbformat.write(nb_sol, f)
                        lab_files.append((nb_s_path, f"Lab_{idx}_{clean_t}_Solution.ipynb"))
                    else:
                        # Conceptual markdown
                        md_content = f"# Lab {idx}: {lab_data.get('lab_title', lab_data.get('title', 'Conceptual Lab'))}\n\n"
                        md_content += f"## Objectives\n{lab_data.get('objectives', '')}\n\n"
                        md_content += f"## Task Description\n{lab_data.get('task_description', '')}\n\n"
                        md_content += f"## Questions\n"
                        for q in lab_data.get("questions", []):
                            md_content += f"- **Question:** {q.get('question', '')}\n"
                            md_content += f"  - **Answer Key:** {q.get('answer', '')}\n\n"
                        
                        md_path = os.path.join(temp_dir, f"Lab_{idx}_{clean_t}.md")
                        with open(md_path, "w", encoding="utf-8") as f:
                            f.write(md_content)
                        lab_files.append((md_path, f"Lab_{idx}_{clean_t}.md"))
                except Exception as lab_err:
                    print(f"Individual lab export skipped: {lab_err}")
                    
        # Create Zip file in memory
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for file_path, zip_path in lab_files:
                if os.path.exists(file_path):
                    zf.write(file_path, zip_path)
                    
        # Cleanup temp files
        for file_path, _ in lab_files:
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except Exception:
                pass
                
        zip_buffer.seek(0)
        zip_path = os.path.join(temp_dir, f"{title.replace(' ', '_')}_labs.zip")
        with open(zip_path, "wb") as f:
            f.write(zip_buffer.getvalue())
            
        return FileResponse(
            path=zip_path,
            filename=f"{title.replace(' ', '_')}_labs.zip",
            media_type="application/zip"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/export_quiz")
async def api_export_quiz(payload: dict = Body(...)):
    title = payload.get("title", "AI_Workshop")
    quiz = payload.get("quiz")
    
    try:
        temp_dir = tempfile.gettempdir()
        
        # 1) DOCX
        docx_temp_path = os.path.join(temp_dir, f"{title.replace(' ', '_')}_quiz.docx")
        has_docx = False
        quiz_data = quiz.get("quiz", quiz) if quiz else None
        if quiz_data and quiz_data.get("questions"):
            try:
                save_quiz_docx(quiz, docx_temp_path)
                has_docx = True
            except Exception as docx_err:
                print(f"DOCX skipped: {docx_err}")
            
        # 2) Kahoot XLSX
        kahoot_temp_path = os.path.join(temp_dir, f"{title.replace(' ', '_')}_kahoot.xlsx")
        kahoot_template = os.path.join(backend_dir, "KahootQuizTemplate.xlsx")
        has_kahoot = False
        if has_docx and os.path.exists(kahoot_template):
            try:
                export_quiz_to_kahoot_xlsx(quiz, kahoot_template, kahoot_temp_path)
                has_kahoot = True
            except Exception as kahoot_err:
                print(f"Kahoot skipped: {kahoot_err}")
            
        # Create Zip file in memory
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            if has_docx and os.path.exists(docx_temp_path):
                zf.write(docx_temp_path, f"{title.replace(' ', '_')}_quiz.docx")
            if has_kahoot and os.path.exists(kahoot_temp_path):
                zf.write(kahoot_temp_path, f"{title.replace(' ', '_')}_kahoot.xlsx")
                
        # Cleanup temp files
        for p in [docx_temp_path, kahoot_temp_path]:
            try:
                if os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass
                
        zip_buffer.seek(0)
        zip_path = os.path.join(temp_dir, f"{title.replace(' ', '_')}_quiz.zip")
        with open(zip_path, "wb") as f:
            f.write(zip_buffer.getvalue())
            
        return FileResponse(
            path=zip_path,
            filename=f"{title.replace(' ', '_')}_quiz.zip",
            media_type="application/zip"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# Import workshop_db
try:
    import workshop_db
except ImportError:
    workshop_db = None

@app.get("/api/workshops")
async def api_list_workshops():
    if not workshop_db:
        return JSONResponse(content=[])
    try:
        workshops = workshop_db.list_workshops()
        return JSONResponse(content=workshops)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/workshops/save")
async def api_save_workshop(payload: dict = Body(...)):
    if not workshop_db:
        raise HTTPException(status_code=500, detail="workshop_db module not available")
    
    title = payload.get("title", "")
    audience = payload.get("audience", "")
    age = payload.get("age", "")
    duration = payload.get("duration", "")
    plan = payload.get("plan", {})
    content = payload.get("content", {})
    labs = payload.get("labs")
    quiz = payload.get("quiz", {})
    
    if not title:
        raise HTTPException(status_code=400, detail="Missing workshop title")
        
    try:
        new_id = workshop_db.save_workshop(
            title=title,
            audience=audience,
            age=age,
            duration=duration,
            plan=plan,
            content=content,
            labs=labs,
            quiz=quiz
        )
        return JSONResponse(content={"id": new_id, "status": "success"})
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/workshops/{workshop_id}")
async def api_load_workshop(workshop_id: int):
    if not workshop_db:
        raise HTTPException(status_code=500, detail="workshop_db module not available")
    try:
        workshop = workshop_db.load_workshop(workshop_id)
        return JSONResponse(content=workshop)
    except ValueError as val_err:
        raise HTTPException(status_code=404, detail=str(val_err))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/workshops/{workshop_id}")
async def api_delete_workshop(workshop_id: int):
    if not workshop_db:
        raise HTTPException(status_code=500, detail="workshop_db module not available")
    try:
        workshop_db.delete_workshop(workshop_id)
        return JSONResponse(content={"status": "success"})
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# Serve Static files from Frontend folder
app.mount("/", StaticFiles(directory="Frontend", html=True), name="static")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("server:app", host="127.0.0.1", port=8000, reload=True)