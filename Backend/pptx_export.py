import os
import copy
import requests
import tempfile
import threading
import contextlib
from concurrent.futures import ThreadPoolExecutor, as_completed
from urllib.parse import urlparse
from search_tool import search_images
from llm_client import ask_llm_to_verify_image
from config import IMAGEN_ASPECT_RATIO, IMAGEN_TIMEOUT_SECONDS  # kept for when Imagen access is confirmed available again
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement


def set_shape_dashed_border(shape, color_rgb):
    """Set shape line to a dashed styled line."""
    ln = shape.line
    ln.color.rgb = color_rgb
    ln.width = Pt(1.5)
    try:
        prstDash = OxmlElement('a:prstDash')
        prstDash.set('val', 'dash')
        ln._line.get_or_add_lnPr().append(prstDash)
    except Exception:
        pass  # Fallback to solid if XML parsing has env issues


# Simple keyword -> icon mapping so a placeholder reads as an intentional
# icon accent instead of debug text like "Brain icon" or "[IMAGE CONCEPT: ...]".
# Not a substitute for real generated images — see note in the handoff below.
ICON_KEYWORDS = [
    (["brain", "reasoning", "cognit"], "🧠"),
    (["wrench", "tool"], "🔧"),
    (["database", "vector", "storage", "ram", "memory"], "🗄️"),
    (["diagram", "flow", "architecture", "process"], "📊"),
    (["timeline", "history", "era", "milestone"], "🕰️"),
    (["robot", "mascot", "agent"], "🤖"),
    (["chat", "conversation", "dialogue"], "💬"),
    (["security", "lock", "auth"], "🔒"),
    (["cloud", "deploy", "scal"], "☁️"),
    (["code", "function", "api", "python"], "💻"),
    (["chart", "graph", "stat"], "📈"),
    (["book", "doc", "knowledge"], "📚"),
]
DEFAULT_ICON = "💡"


def pick_icon(concept_text: str) -> str:
    text = (concept_text or "").lower()
    for keywords, emoji in ICON_KEYWORDS:
        if any(k in text for k in keywords):
            return emoji
    return DEFAULT_ICON

import re

IMAGE_MAP = {
    "security": "https://images.unsplash.com/photo-1510511459019-5dda7724fd87?w=600&q=80",
    "lock": "https://images.unsplash.com/photo-1510511459019-5dda7724fd87?w=600&q=80",
    "padlock": "https://images.unsplash.com/photo-1510511459019-5dda7724fd87?w=600&q=80",
    "crypt": "https://images.unsplash.com/photo-1526374965328-7f61d4dc18c5?w=600&q=80",
    "api": "https://images.unsplash.com/photo-1558494949-ef010cbdcc31?w=600&q=80",
    "server": "https://images.unsplash.com/photo-1558494949-ef010cbdcc31?w=600&q=80",
    "database": "https://images.unsplash.com/photo-1448375240586-882707db888b?w=600&q=80",
    "code": "https://images.unsplash.com/photo-1542831371-29b0f74f9713?w=600&q=80",
    "program": "https://images.unsplash.com/photo-1542831371-29b0f74f9713?w=600&q=80",
    "python": "https://images.unsplash.com/photo-1542831371-29b0f74f9713?w=600&q=80",
    "ai": "https://images.unsplash.com/photo-1677442136019-21780efad99a?w=600&q=80",
    "robot": "https://images.unsplash.com/photo-1485827404703-89b55fcc595e?w=600&q=80",
    "agent": "https://images.unsplash.com/photo-1485827404703-89b55fcc595e?w=600&q=80",
    "poetry": "https://images.unsplash.com/photo-1512820790803-83ca734da794?w=600&q=80",
    "book": "https://images.unsplash.com/photo-1512820790803-83ca734da794?w=600&q=80",
    "write": "https://images.unsplash.com/photo-1455390582262-044cdead277a?w=600&q=80",
    "timeline": "https://images.unsplash.com/photo-1508962914676-134849a727f0?w=600&q=80",
    "history": "https://images.unsplash.com/photo-1508962914676-134849a727f0?w=600&q=80",
    "diagram": "https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=600&q=80",
    "chart": "https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=600&q=80",
    "network": "https://images.unsplash.com/photo-1544197150-b99a580bb7a8?w=600&q=80",
}

USED_KEYWORDS = set()

def extract_keywords(concept_text):
    if not concept_text:
        return "technology"
        
    text = concept_text.lower()
    
    # 1. Direct visual noun mapping to match concrete physical objects
    visual_nouns = {
        "padlock": "lock",
        "lock": "lock",
        "key": "key",
        "server": "server",
        "database": "database",
        "folder": "folder",
        "document": "document",
        "file": "file",
        "code": "programming",
        "program": "programming",
        "python": "coding",
        "robot": "robot",
        "agent": "ai",
        "ai": "technology",
        "poetry": "poetry",
        "poem": "writing",
        "book": "books",
        "write": "writing",
        "classroom": "classroom",
        "school": "school",
        "kids": "children",
        "kid": "children",
        "children": "children",
        "clock": "clock",
        "timeline": "timeline",
        "time": "clock",
        "break": "coffee",
        "coffee": "coffee",
        "discuss": "meeting",
        "question": "question",
        "mascot": "robot",
        "security": "security",
        "crypt": "cryptography",
        "cyber": "security",
        "network": "network",
        "cloud": "cloud",
        "chart": "chart",
        "diagram": "diagram",
        "graph": "graph",
    }
    
    # If the visual concept contains a concrete visual noun, use it directly to ensure perfect image relevance
    for key, search_term in visual_nouns.items():
        if key in text:
            return search_term
            
    # 2. Fallback: extract the first 2 significant words
    for char in [".", ",", "!", "?", ";", ":", "-", "_", "(", ")", "[", "]", "{", "}"]:
        text = text.replace(char, " ")
        
    filler = ["icon", "representing", "concept", "illustration", "showing", "graphic", "picture", "image", "a", "an", "the", "of", "for", "with", "over", "under", "in", "on", "at", "by"]
    words = text.split()
    filtered = [w for w in words if w not in filler and len(w) > 2]
    
    if filtered:
        return ",".join(filtered[:2])
        
    return "technology"

def collect_image_concepts(slides: list[dict]) -> set:
    """
    Walks every slide's blocks and pulls out every distinct image concept
    text that will need a real image fetched for it — image_placeholder
    blocks directly, plus nested image_placeholder fields inside
    columns_3 columns. Used to fetch every needed image UP FRONT, in
    parallel, instead of one at a time while slides are being built.
    """
    concepts = set()
    for slide in slides:
        for block in slide.get("blocks", []):
            if block.get("type") == "image_placeholder" and block.get("text"):
                concepts.add(block["text"])
            if block.get("type") == "columns_3":
                for col in block.get("columns", []):
                    if col.get("image_placeholder"):
                        concepts.add(col["image_placeholder"])
    return concepts


def prefetch_images(concepts: set, max_workers: int = 5) -> dict:
    """
    Fetches every concept in `concepts` in parallel (bounded to
    max_workers at once) instead of sequentially — this is the single
    biggest lever for PPTX export speed, since a deck with ~15-20 image
    slots previously fetched each one strictly one after another.

    Returns {concept_text: local_file_path_or_None}, meant to be passed
    into export_to_pptx as `prefetched_images` and forwarded down to
    every download_image_for_concept call, which checks this cache
    before ever doing a live search.

    Thread safety: used_image_ids (the "don't repeat a photo in this
    deck" set) is shared and mutated from multiple worker threads here,
    so it's guarded by a lock — the ONLY part of download_image_for_concept
    that touches shared state across concurrent calls.
    """
    if not concepts:
        return {}

    used_image_ids = set()
    lock = threading.Lock()
    results = {}

    def _fetch_one(concept_text: str):
        return concept_text, download_image_for_concept(concept_text, used_image_ids=used_image_ids, lock=lock)

    with ThreadPoolExecutor(max_workers=min(max_workers, len(concepts))) as executor:
        futures = [executor.submit(_fetch_one, c) for c in concepts]
        for future in as_completed(futures):
            concept_text, path = future.result()
            results[concept_text] = path
    return results


def download_image_for_concept(concept_text, used_image_ids=None, lock=None):
    """
    used_image_ids: an optional set the CALLER maintains across one whole
    presentation export. When any source returns an image, its
    identifying key gets added here — if the SAME key would be picked
    again for a later slide, we skip it and try the next option, so one
    deck never shows the same image twice.

    lock: an optional threading.Lock, required when this function may be
    called concurrently from multiple threads (see prefetch_images) —
    guards every read/write of used_image_ids so two threads can't both
    "claim" the same image at once. Safe to omit when called from a
    single thread (used_image_ids access is then unguarded, same as
    before this parameter existed).

    Tries, in order:
      1. Tavily image search — REAL images tied to actual matching web
         pages, free (same API key already used for text search). This
         is the primary source: relevant and reliable. WEBP results are
         skipped (python-pptx can't embed WEBP at all).
      2. A small fixed local image pool (IMAGE_MAP), for when nothing
         from search fits.
      3. None — falls through to the dashed-box-with-icon placeholder
         (see add_visual_placeholder) rather than showing a wrong image.

    Imagen 3 generation and Loremflickr's random-tag lookup used to sit
    in this chain and were both removed — Imagen 3 reliably 404s on this
    account/API version, and Loremflickr repeatedly returned images with
    no real connection to the concept. See the code below for specifics.
    """
    if not concept_text:
        return None
    if used_image_ids is None:
        used_image_ids = set()
    guard = lock if lock is not None else contextlib.nullcontext()

    # 1. Tavily image search — try this FIRST, before the paid/unreliable
    # options below. Real, contextually relevant images, and usually
    # fast — this also cuts typical export time, since we're no longer
    # waiting out an Imagen failure (or a billing rejection) on every
    # single slide before getting a usable image.
    try:
        image_results = search_images(
            f"{concept_text} technical diagram illustration", max_results=3
        )
        for result in image_results:
            url = result["url"]
            with guard:
                if url in used_image_ids:
                    continue
            try:
                r = requests.get(url, timeout=4)
                content_type = r.headers.get("Content-Type", "")
                # python-pptx's add_picture only accepts BMP/GIF/JPEG/PNG/
                # TIFF/WMF — WEBP (increasingly common on modern sites) is
                # NOT supported and throws an unhandled exception that
                # crashes the whole export. Skip WEBP results outright and
                # try the next search candidate instead.
                if r.status_code == 200 and "image" in content_type and "webp" not in content_type.lower():
                    # Verify the image ACTUALLY shows this concept before
                    # accepting it — search relevance is text-based (the
                    # page around the image matched), not a guarantee the
                    # image itself is relevant. Bounded: this runs at most
                    # once per candidate, and there are at most 3
                    # candidates per concept (max_results=3 above), so this
                    # can never turn into an open-ended search.
                    if not ask_llm_to_verify_image(r.content, concept_text):
                        continue  # this candidate isn't actually relevant — try the next
                    fd, path = tempfile.mkstemp(suffix=".jpg")
                    with os.fdopen(fd, 'wb') as tmp:
                        tmp.write(r.content)
                    with guard:
                        used_image_ids.add(url)
                    return path
            except Exception:
                continue  # this particular image URL didn't work, try the next
    except Exception as e:
        print(f"Tavily image search failed, falling back to other sources: {e}")

    # Imagen 3 used to sit here as a second-tier source. Removed: confirmed
    # via a real error from Google's API that "imagen-3.0-generate-002 is
    # not found for API version v1beta, or is not supported for predict"
    # for this account — every single call was a guaranteed, wasted 404,
    # not just an occasional billing failure. If Imagen access is ever
    # enabled for this project later, this is the place to re-add it —
    # but don't re-add it speculatively; confirm the model is reachable
    # first (see docs.google's model list for the exact current name).

    # Loremflickr (a "random photo tagged roughly like this word" service)
    # used to sit here as a fallback. Removed: it repeatedly returned
    # actively WRONG images (e.g. an unrelated statue photo, a financial
    # news thumbnail for "guardrail") rather than no image at all — which
    # is worse than the honest dashed-box-with-caption placeholder below.
    # A curated, keyword-matched url (IMAGE_MAP) is a safer fallback than
    # a blind tag lottery.
    kw = extract_keywords(concept_text)

    # Fallback to local map. Several DIFFERENT keywords map to the SAME
    # url BY DESIGN (e.g. "security"/"lock"/"padlock" all point at one
    # photo) — so two different slide concepts ("security" on one slide,
    # "lock" on another) can each correctly match the map and still land
    # on an identical picture. To avoid that: try the keyword's own match
    # first; if THAT exact url was already used earlier in this export,
    # fall back to ANY other still-unused url in the whole map — a
    # same-topic-adjacent photo beats a guaranteed duplicate. Only reuse
    # a url if the entire pool has been used.
    matched_url = None
    for k, val in IMAGE_MAP.items():
        if k in kw:
            matched_url = val
            break

    with guard:
        if matched_url and matched_url not in used_image_ids:
            ordered_candidates = [matched_url]
        else:
            ordered_candidates = [v for v in IMAGE_MAP.values() if v not in used_image_ids]
        # Always end with every remaining map value (even ones already
        # used) as a last resort, so a dead/unreachable URL earlier in
        # the list doesn't leave this concept with no image at all.
        ordered_candidates += [v for v in IMAGE_MAP.values() if v not in ordered_candidates]

    for url_fallback in ordered_candidates:
        try:
            r = requests.get(url_fallback, timeout=3)
            if r.status_code == 200:
                fd, path = tempfile.mkstemp(suffix=".jpg")
                with os.fdopen(fd, 'wb') as tmp:
                    tmp.write(r.content)
                with guard:
                    used_image_ids.add(url_fallback)
                return path
        except Exception:
            continue  # this particular fallback URL didn't work, try the next

    return None




def add_visual_placeholder(slide, left, top, width, height, concept_text, slide_style, cat_color, temp_files=None, include_images=True, used_image_ids=None, prefetched_images=None):
    # Use the pre-fetched image (see prefetch_images) if this concept was
    # already looked up in the parallel prefetch pass — avoids a second,
    # redundant live search here. Falls back to a fresh live search only
    # if prefetched_images wasn't provided or doesn't have this concept
    # (e.g. the small set of hardcoded structural concepts like the
    # roadmap/thank-you slide's own decorative image, which aren't worth
    # including in the bulk prefetch since there's only ever one or two
    # of them per deck).
    if include_images and prefetched_images is not None and concept_text in prefetched_images:
        img_path = prefetched_images[concept_text]
    else:
        img_path = download_image_for_concept(concept_text, used_image_ids=used_image_ids) if include_images else None
    if img_path and os.path.exists(img_path):
        try:
            pic = slide.shapes.add_picture(img_path, left, top, width, height)
            if temp_files is not None:
                temp_files.append(img_path)
            else:
                try:
                    os.remove(img_path)
                except Exception:
                    pass
            return pic
        except Exception as pic_err:
            print(f"Failed to add downloaded picture: {pic_err}")

    # Fallback to dashed box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.06
    box.fill.solid()
    if slide_style == "Bold & Impactful":
        box.fill.fore_color.rgb = RGBColor(34, 40, 66)
    else:
        box.fill.fore_color.rgb = RGBColor(238, 244, 250)
    set_shape_dashed_border(box, cat_color)

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    icon_p = tf.paragraphs[0]
    icon_p.alignment = PP_ALIGN.CENTER
    icon_p.text = pick_icon(concept_text)
    icon_p.font.size = Pt(40)

    if concept_text:
        caption = concept_text if len(concept_text) <= 220 else concept_text[:217] + "..."
        cap_p = tf.add_paragraph()
        cap_p.alignment = PP_ALIGN.CENTER
        cap_p.text = caption
        cap_p.font.name = "Calibri"
        cap_p.font.size = Pt(9)
        cap_p.font.italic = True
        cap_p.font.color.rgb = RGBColor(113, 128, 150)
        cap_p.space_before = Pt(6)

    return box


def apply_slide_background_and_style(slide, style):
    """Sets background and returns (title_color, cat_color, body_color, title_font, body_font)."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    rect.line.fill.background()

    bg_color = RGBColor(255, 255, 255)
    title_color = RGBColor(0, 13, 32)
    cat_color = RGBColor(102, 126, 234)
    body_color = RGBColor(74, 85, 104)
    title_font = "Arial"
    body_font = "Calibri"

    if style == "Bold & Impactful":
        bg_color = RGBColor(26, 26, 46)
        title_color = RGBColor(255, 255, 255)
        cat_color = RGBColor(121, 40, 202)
        body_color = RGBColor(226, 232, 240)
        title_font = "Arial Black"
        body_font = "Arial"
    elif style == "Visual & Diagram-heavy":
        bg_color = RGBColor(244, 248, 252)
        title_color = RGBColor(0, 13, 32)
        cat_color = RGBColor(102, 126, 234)
        body_color = RGBColor(74, 85, 104)
        title_font = "Segoe UI"
        body_font = "Segoe UI"
    elif style == "Data & Research":
        bg_color = RGBColor(255, 255, 255)
        title_color = RGBColor(17, 17, 17)
        cat_color = RGBColor(60, 60, 60)
        body_color = RGBColor(40, 40, 40)
        title_font = "Georgia"
        body_font = "Georgia"
        border_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(12.533), Inches(6.7)
        )
        border_rect.fill.background()
        border_rect.line.color.rgb = RGBColor(200, 200, 200)
        border_rect.line.width = Pt(1)
    elif style == "Interactive & Workshop":
        bg_color = RGBColor(250, 248, 245)
        title_color = RGBColor(45, 55, 72)
        cat_color = RGBColor(221, 107, 32)
        body_color = RGBColor(74, 85, 104)
        title_font = "Trebuchet MS"
        body_font = "Trebuchet MS"
    else:  # Clean & Minimal
        bg_color = RGBColor(255, 255, 255)
        title_color = RGBColor(17, 17, 17)
        cat_color = RGBColor(85, 85, 85)
        body_color = RGBColor(60, 60, 60)
        # NOTE: previously drew a divider line under the title here.
        # Removed — accent lines under titles are a well-known AI-slide
        # tell. Whitespace + the background color already differentiate
        # this style; no extra element needed.

    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color

    try:
        slide.shapes._spTree.remove(rect._element)
        slide.shapes._spTree.insert(2, rect._element)
    except Exception:
        pass

    return title_color, cat_color, body_color, title_font, body_font


def add_category_header(slide, section_name, cat_color, font_name):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = str(section_name).upper()
    p.font.name = font_name
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = cat_color


def add_role_badge(slide, role_label, cat_color, font_name, title_color):
    """
    Small pill-style label (e.g. "LAB", "EXPLAIN", "BREAK") in the top
    right corner, tinted with the section's own accent color — a quick,
    low-cost way to make a slide's type visible at a glance without a
    full redesign, and without anything that risks looking heavy or
    slowing down export (this is a single small shape + one line of text,
    no images, no extra API calls).
    """
    if not role_label:
        return
    badge_width = Inches(1.3)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(13.333) - badge_width - Inches(0.5), Inches(0.35), badge_width, Inches(0.35)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = cat_color
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(role_label).upper()
    p.font.name = font_name
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def add_page_number(slide, page_num, total_pages, body_color, font_name):
    """
    Small "N / total" label, bottom-right — a standard, low-cost polish
    touch (no images, no extra API calls, pure static text) that reads as
    professional without adding anything that could slow export down.
    """
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.3))
    tf = txBox.text_frame
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = f"{page_num} / {total_pages}"
    p.font.name = font_name
    p.font.size = Pt(9)
    p.font.color.rgb = body_color


def format_source_footer(sources: list) -> str:
    """Domain only, not the full raw URL — reads as a real citation, not a debug dump."""
    if not sources:
        return ""
    # Guard against string-type source entries from LLM
    for source in sources:
        if not isinstance(source, dict):
            continue
        first_url = source.get("url", "")
        if first_url:
            domain = urlparse(first_url).netloc.replace("www.", "")
            suffix = f" +{len(sources) - 1} more" if len(sources) > 1 else ""
            return f"SOURCE  {domain}{suffix}"
    return ""


def extract_template_styles(prs):
    # Defaults
    styles = {
        "title_font": "Arial",
        "title_color": RGBColor(0, 13, 32),
        "body_font": "Calibri",
        "body_color": RGBColor(74, 85, 104),
        "cat_color": RGBColor(102, 126, 234),
    }

    # 1. Check existing slides (even if empty, check format properties)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                for p in tf.paragraphs:
                    is_title = (shape == slide.shapes.title or 
                                (p.font.size and p.font.size > Pt(24)) or 
                                (shape.is_placeholder and shape.placeholder_format.type == 1))
                    if is_title:
                        if p.font.name:
                            styles["title_font"] = p.font.name
                        if p.font.color and p.font.color.type == 1:
                            styles["title_color"] = p.font.color.rgb
                    else:
                        if p.font.name:
                            styles["body_font"] = p.font.name
                        if p.font.color and p.font.color.type == 1:
                            styles["body_color"] = p.font.color.rgb

    # 2. Check layouts
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                tf = shape.text_frame
                p = tf.paragraphs[0] if tf.paragraphs else None
                if p:
                    if ph_type == 1: # Title
                        if p.font.name:
                            styles["title_font"] = p.font.name
                        if p.font.color and p.font.color.type == 1:
                            styles["title_color"] = p.font.color.rgb
                    elif ph_type in [2, 7]: # Body / Content
                        if p.font.name:
                            styles["body_font"] = p.font.name
                        if p.font.color and p.font.color.type == 1:
                            styles["body_color"] = p.font.color.rgb

    styles["cat_color"] = styles["title_color"]
    return styles


def find_content_layout(prs):
    # Stencils will be located directly in the slides, so layouts are fallback only
    if len(prs.slide_layouts) > 1:
        return prs.slide_layouts[1]
    return prs.slide_layouts[0]


def _copy_background(src_slide, dst_slide):
    """
    Copy only the background fill XML from src to dst without touching
    any image relationships (which would corrupt the target file).
    We copy the <p:bg> element only, which carries solid fills, gradients,
    and pattern fills safely because it has no rId references.
    """
    try:
        src_bg = src_slide.element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
        if src_bg is None:
            return
        # Check if the background references an image (r:embed / r:link) — skip if so
        ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        if src_bg.find(f'.//{{{ns}}}embed') is not None or src_bg.find(f'.//{{{ns}}}link') is not None:
            return
        import copy as _copy
        # Remove any existing bg element
        existing = dst_slide.element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
        if existing is not None:
            existing.getparent().remove(existing)
        # Insert the copied bg at index 0 of the slide element (which is required by the OpenXML schema)
        dst_slide.element.insert(0, _copy.deepcopy(src_bg))
    except Exception as bg_err:
        print(f"Background copy skipped: {bg_err}")


def export_to_pptx(title, slides, output_path="workshop.pptx", template_path=None, include_images=True):
    """Export slides to PowerPoint. When a template is provided, its slide master
    (theme, layouts, fonts, background colors) is preserved. Slides are built
    programmatically — no shape cloning which causes PPTX corruption."""

    is_template = False
    styles = None

    # Tracks every image (Imagen call, loremflickr download, or IMAGE_MAP
    # url) used so far in THIS export — passed down to every
    # add_visual_placeholder call below so the same photo never appears
    # twice across one presentation.
    used_image_ids = set()

    # Fetch every needed image UP FRONT, in parallel, instead of one at a
    # time while slides are being built — this is what previously made
    # exporting a 15-20 image deck slow (each fetch waited for the last
    # one to finish before starting). See prefetch_images' docstring.
    prefetched_images = prefetch_images(collect_image_concepts(slides)) if include_images else {}

    if template_path and os.path.exists(template_path):
        try:
            # Load the template to extract styles FIRST (before clearing slides)
            stencil_prs = Presentation(template_path)
            styles = extract_template_styles(stencil_prs)
            # Keep a reference to a content slide background for copying
            bg_source_slide = stencil_prs.slides[1] if len(stencil_prs.slides) > 1 else (stencil_prs.slides[0] if stencil_prs.slides else None)
            cover_bg_slide = stencil_prs.slides[0] if stencil_prs.slides else None

            # Now load a fresh copy and clear all slides — theme/master is preserved
            prs = Presentation(template_path)
            id_list = prs.slides._sldIdLst
            for i in range(len(id_list) - 1, -1, -1):
                slide_id = id_list[i]
                prs.part.drop_rel(slide_id.rId)
                del id_list[i]

            is_template = True
        except Exception as e:
            print(f"Failed to load template: {e}. Falling back to default.")
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # Pick the blank layout (index 6 is blank in most templates; fall back to last)
    def blank_layout():
        try:
            return prs.slide_layouts[6]
        except IndexError:
            return prs.slide_layouts[-1]

    logo_path = os.path.join(os.path.dirname(__file__), "logo.png")

    def get_colors(slide_style="Clean & Minimal"):
        if is_template and styles:
            return (
                styles["title_color"],
                styles["cat_color"],
                styles["body_color"],
                styles["title_font"],
                styles["body_font"],
            )
        return apply_slide_background_and_style.__wrapped__(slide_style) if hasattr(apply_slide_background_and_style, '__wrapped__') else (None, None, None, None, None)

    def setup_slide(slide, slide_style, use_bg_source=None):
        """Apply background and return color/font tuple."""
        if is_template and styles:
            # Apply only background color from master (safe, no rId references)
            if use_bg_source:
                _copy_background(use_bg_source, slide)
            tc = styles["title_color"]
            cc = styles["cat_color"]
            bc = styles["body_color"]
            tf = styles["title_font"]
            bf = styles["body_font"]
        else:
            tc, cc, bc, tf, bf = apply_slide_background_and_style(slide, slide_style)
        # Add logo
        if os.path.exists(logo_path):
            try:
                slide.shapes.add_picture(logo_path, Inches(11.2), Inches(0.3), width=Inches(1.6))
            except Exception:
                pass
        return tc, cc, bc, tf, bf

    def add_title_box(slide, text, tc, tf, top=Inches(0.7), left=Inches(0.8), width=Inches(9.8)):
        font_size = Pt(28) if len(text) > 30 else Pt(36)
        tb = slide.shapes.add_textbox(left, top, width, Inches(1.4))
        tframe = tb.text_frame
        tframe.word_wrap = True
        tframe.margin_left = tframe.margin_top = tframe.margin_bottom = tframe.margin_right = 0
        p = tframe.paragraphs[0]
        p.text = text
        p.font.name = tf
        p.font.size = font_size
        p.font.bold = True
        p.font.color.rgb = tc
        return tb

    # ── TITLE SLIDE ─────────────────────────────────────────────────────────
    title_slide = prs.slides.add_slide(blank_layout())
    first_style = slides[0].get('slide_style', 'Clean & Minimal') if slides else 'Clean & Minimal'
    tc, cc, bc, tf, bf = setup_slide(title_slide, first_style, use_bg_source=cover_bg_slide if is_template else None)

    # Big centered title. Font size AND the subtitle's vertical position
    # both scale with title length — a long title (this one wraps to 4
    # lines at 44pt) was previously overlapping the fixed-position
    # "WORKSHOP PRESENTATION" subtitle below it, since the title textbox
    # doesn't auto-shrink and nothing accounted for how many lines a long
    # title would actually take.
    title_len = len(title)
    if title_len > 90:
        title_font_size = Pt(28)
        title_box_height = Inches(3.2)
        subtitle_top = Inches(5.5)
    elif title_len > 60:
        title_font_size = Pt(32)
        title_box_height = Inches(2.8)
        subtitle_top = Inches(5.1)
    elif title_len > 30:
        title_font_size = Pt(38)
        title_box_height = Inches(2.5)
        subtitle_top = Inches(4.7)
    else:
        title_font_size = Pt(44)
        title_box_height = Inches(2.2)
        subtitle_top = Inches(4.3)

    tb = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), title_box_height)
    tframe = tb.text_frame
    tframe.word_wrap = True
    p = tframe.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = title.upper()
    p.font.name = tf
    p.font.size = title_font_size
    p.font.bold = True
    p.font.color.rgb = tc

    sb = title_slide.shapes.add_textbox(Inches(1.0), subtitle_top, Inches(11.333), Inches(0.8))
    sf = sb.text_frame
    sf.word_wrap = True
    sp = sf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sp.text = "WORKSHOP PRESENTATION"
    sp.font.name = bf
    sp.font.size = Pt(20)
    sp.font.color.rgb = bc

    temp_images = []

    # ── ROADMAP SLIDE ────────────────────────────────────────────────────────
    roadmap_slide = prs.slides.add_slide(blank_layout())
    tc, cc, bc, tf, bf = setup_slide(roadmap_slide, first_style, use_bg_source=bg_source_slide if is_template else None)

    roadmap_points = []
    seen = set()
    for s in slides:
        sec = s.get('section', '')
        if sec and sec not in seen:
            seen.add(sec)
            roadmap_points.append(sec)
    if not roadmap_points:
        for s in slides:
            t = s.get('slide_title', '')
            if t and t not in seen:
                seen.add(t)
                roadmap_points.append(t)
    roadmap_points = roadmap_points[:8]

    rm_title_text = "ورقة العمل / خطة الورشة" if any(ord(c) > 127 for c in title) else "COURSE ROADMAP"
    add_title_box(roadmap_slide, rm_title_text, tc, tf)
    add_category_header(roadmap_slide, "Overview", cc, tf)

    list_box = roadmap_slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(7.2), Inches(4.5))
    ltf = list_box.text_frame
    ltf.word_wrap = True
    first_p = True
    for idx, pt in enumerate(roadmap_points, 1):
        p = ltf.paragraphs[0] if first_p else ltf.add_paragraph()
        first_p = False
        p.text = f"  {idx}.  {pt}"
        p.font.name = tf
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = tc
        p.space_before = Pt(10)
        p.space_after = Pt(4)

    add_visual_placeholder(
        roadmap_slide, Inches(8.3), Inches(2.2), Inches(4.2), Inches(4.3),
        "roadmap overview agenda", first_style, cc, temp_files=temp_images, include_images=include_images, used_image_ids=used_image_ids, prefetched_images=prefetched_images
    )

    # ── CONTENT SLIDES ───────────────────────────────────────────────────────
    for slide_data in slides:
        c_type = slide_data.get('content_type', 'content_slide')
        slide_style = slide_data.get('slide_style', 'Clean & Minimal')

        slide = prs.slides.add_slide(blank_layout())
        tc, cc, bc, tf, bf = setup_slide(slide, slide_style, use_bg_source=bg_source_slide if is_template else None)

        section_name = slide_data.get('section', 'Workshop')
        add_category_header(slide, section_name, cc, tf)

        title_text = slide_data.get('slide_title', '').upper()
        add_title_box(slide, title_text, tc, tf)

        content_top = Inches(2.3)

        # ── Roadmap slide ────────────────────────────────────────────────────
        if c_type == 'roadmap_slide':
            roadmap_items = []
            for block in slide_data.get('blocks', []):
                if block.get('type') == 'roadmap':
                    roadmap_items = block.get('items', [])

            list_box = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(7.0), Inches(4.0))
            ltf = list_box.text_frame
            ltf.word_wrap = True
            ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ltf.margin_left = ltf.margin_top = ltf.margin_bottom = ltf.margin_right = 0
            first_p = True
            for item in roadmap_items:
                p_title = ltf.paragraphs[0] if first_p else ltf.add_paragraph()
                first_p = False
                p_title.text = f"{item.get('title', '')}"
                p_title.font.name = tf
                p_title.font.size = Pt(20)
                p_title.font.bold = True
                p_title.font.color.rgb = tc
                p_title.space_before = Pt(14)
                p_desc = ltf.add_paragraph()
                p_desc.text = f"  ● {item.get('detail', '')}"
                p_desc.font.name = bf
                p_desc.font.size = Pt(15)
                p_desc.font.color.rgb = bc
                p_desc.space_after = Pt(10)

            add_visual_placeholder(
                slide, Inches(8.5), content_top, Inches(4.0), Inches(4.0),
                "Workshop mascot illustration", slide_style, cc, temp_files=temp_images, include_images=include_images, used_image_ids=used_image_ids, prefetched_images=prefetched_images
            )

        # ── 3-columns slide ──────────────────────────────────────────────────
        elif c_type == 'columns_3_slide':
            subtitle_text = ""
            cols_items = []
            for block in slide_data.get('blocks', []):
                if block.get('type') == 'paragraph':
                    subtitle_text = block.get('text', '')
                elif block.get('type') == 'columns_3':
                    cols_items = block.get('columns', [])

            if subtitle_text:
                sub_box = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(11.7), Inches(0.6))
                stf = sub_box.text_frame
                stf.word_wrap = True
                stf.margin_left = stf.margin_top = stf.margin_bottom = stf.margin_right = 0
                sp2 = stf.paragraphs[0]
                sp2.text = subtitle_text
                sp2.font.name = bf
                sp2.font.size = Pt(18)
                sp2.font.color.rgb = bc

            col_width = Inches(3.6)
            col_gap = Inches(0.4)
            start_left = Inches(0.8)
            col_top = content_top + (Inches(0.8) if subtitle_text else Inches(0.2))

            for idx, col_data in enumerate(cols_items[:3]):
                curr_left = start_left + idx * (col_width + col_gap)
                col_box = slide.shapes.add_textbox(curr_left, col_top, col_width, Inches(1.4))
                ctf = col_box.text_frame
                ctf.word_wrap = True
                ctf.margin_left = ctf.margin_top = ctf.margin_bottom = ctf.margin_right = 0
                p_head = ctf.paragraphs[0]
                p_head.text = f"● {col_data.get('heading', '')}"
                p_head.font.name = tf
                p_head.font.size = Pt(18)
                p_head.font.bold = True
                p_head.font.color.rgb = tc
                p_head.space_after = Pt(6)
                p_desc = ctf.add_paragraph()
                p_desc.text = col_data.get('text', '')
                p_desc.font.name = bf
                p_desc.font.size = Pt(13)
                p_desc.font.color.rgb = bc
                p_desc.space_after = Pt(10)
                add_visual_placeholder(
                    slide, curr_left, col_top + Inches(1.4), col_width, Inches(2.2),
                    col_data.get('image_placeholder', ''), slide_style, cc, temp_files=temp_images, include_images=include_images, used_image_ids=used_image_ids, prefetched_images=prefetched_images
                )

        # ── Timeline slide ───────────────────────────────────────────────────
        elif c_type == 'timeline_slide':
            timeline_items = []
            right_image_concept = "Timeline graphic"
            for block in slide_data.get('blocks', []):
                if block.get('type') == 'timeline':
                    timeline_items = block.get('events', [])
                elif block.get('type') == 'image_placeholder':
                    right_image_concept = block.get('text', 'Timeline graphic')

            tl_box = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(7.5), Inches(4.0))
            ttf = tl_box.text_frame
            ttf.word_wrap = True
            ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ttf.margin_left = ttf.margin_top = ttf.margin_bottom = ttf.margin_right = 0
            first_p = True
            for item in timeline_items:
                p_evt = ttf.paragraphs[0] if first_p else ttf.add_paragraph()
                first_p = False
                p_evt.text = f"➔ {item.get('title', '')} · {item.get('date', '')}"
                p_evt.font.name = tf
                p_evt.font.size = Pt(14)
                p_evt.font.bold = True
                p_evt.font.color.rgb = tc
                p_evt.space_before = Pt(8)
                p_desc = ttf.add_paragraph()
                p_desc.text = f"    {item.get('text', '')}"
                p_desc.font.name = bf
                p_desc.font.size = Pt(12)
                p_desc.font.color.rgb = bc
                p_desc.space_after = Pt(6)

            add_visual_placeholder(
                slide, Inches(8.8), content_top + Inches(0.4), Inches(3.8), Inches(3.7),
                right_image_concept, slide_style, cc, temp_files=temp_images, include_images=include_images, used_image_ids=used_image_ids, prefetched_images=prefetched_images
            )

        # ── Standard content slide ───────────────────────────────────────────
        else:
            blocks = slide_data.get('blocks', [])
            has_activity = any(b.get('type') == 'activity' for b in blocks)
            has_image = any(b.get('type') == 'image_placeholder' for b in blocks)
            image_block = next((b for b in blocks if b.get('type') == 'image_placeholder'), None)

            body_height = Inches(2.0) if has_activity else Inches(4.0)
            body_width = Inches(7.0) if has_image else Inches(11.7)

            body_box = slide.shapes.add_textbox(Inches(0.8), content_top, body_width, body_height)
            btf = body_box.text_frame
            btf.word_wrap = True
            btf.margin_left = btf.margin_top = btf.margin_bottom = btf.margin_right = 0

            first_p = True
            for block in blocks:
                b_type = block.get('type')
                text = block.get('text', '')

                if b_type == 'heading':
                    p = btf.paragraphs[0] if first_p else btf.add_paragraph()
                    first_p = False
                    p.text = str(text)
                    p.font.name = tf
                    p.font.bold = True
                    p.font.size = Pt(22)
                    p.font.color.rgb = cc
                    p.space_before = Pt(10)
                    p.space_after = Pt(8)
                elif b_type == 'paragraph':
                    p = btf.paragraphs[0] if first_p else btf.add_paragraph()
                    first_p = False
                    p.text = str(text)
                    p.font.name = bf
                    p.font.size = Pt(16)
                    p.font.color.rgb = bc
                    p.space_after = Pt(10)
                elif b_type == 'bullet_list':
                    for item in block.get('items', []):
                        p = btf.paragraphs[0] if first_p else btf.add_paragraph()
                        first_p = False
                        p.text = str(item)
                        p.level = 1
                        p.font.name = bf
                        p.font.size = Pt(15)
                        p.font.color.rgb = bc
                        p.space_after = Pt(6)
                elif b_type == 'activity':
                    act_shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.733), Inches(1.8)
                    )
                    act_shape.adjustments[0] = 0.04
                    act_shape.fill.solid()
                    if slide_style == 'Bold & Impactful':
                        act_shape.fill.fore_color.rgb = RGBColor(50, 20, 70)
                        act_shape.line.color.rgb = RGBColor(221, 107, 32)
                    else:
                        act_shape.fill.fore_color.rgb = RGBColor(255, 245, 235)
                        act_shape.line.color.rgb = RGBColor(221, 107, 32)
                    atf = act_shape.text_frame
                    atf.word_wrap = True
                    atf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    ap = atf.paragraphs[0]
                    ap.text = f"⚡ ACTIVITY: {text}"
                    ap.font.name = tf
                    ap.font.size = Pt(14)
                    ap.font.bold = True
                    ap.font.color.rgb = RGBColor(255, 255, 255) if slide_style == 'Bold & Impactful' else RGBColor(221, 107, 32)

            if image_block:
                add_visual_placeholder(
                    slide, Inches(8.3), content_top, Inches(4.2), Inches(4.0),
                    image_block.get('text', ''), slide_style, cc, temp_files=temp_images, include_images=include_images, used_image_ids=used_image_ids, prefetched_images=prefetched_images
                )

        # Speaker notes
        notes_slide = slide.notes_slide
        notes_text = "Sources:\n"
        for source in slide_data.get('sources', []):
            if not isinstance(source, dict):
                notes_text += f"- {source}\n"
                continue
            notes_text += f"- {source.get('author', '')} ({source.get('year', '')}): {source.get('title', '')} - {source.get('url', '')}\n"
            takeaway = source.get('takeaway') or source.get('exact_quote')
            if takeaway:
                notes_text += f"  Note: {takeaway}\n"
        notes_slide.notes_text_frame.text = notes_text

        # Source footer
        footer_text = format_source_footer(slide_data.get('sources', []))
        if footer_text:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(10.5), Inches(0.35))
            tf_f = txBox.text_frame
            tf_f.word_wrap = True
            p_f = tf_f.paragraphs[0]
            p_f.text = footer_text
            p_f.font.name = bf
            p_f.font.size = Pt(8)
            p_f.font.bold = True
            p_f.font.color.rgb = RGBColor(160, 174, 192)

    # ── THANK YOU SLIDE ──────────────────────────────────────────────────────
    thank_you_slide = prs.slides.add_slide(blank_layout())
    ty_style = slides[-1].get('slide_style', 'Clean & Minimal') if slides else 'Clean & Minimal'
    tc, cc, bc, tf, bf = setup_slide(thank_you_slide, ty_style, use_bg_source=cover_bg_slide if is_template else None)

    # Big centered "Thank You" text
    ty_box = thank_you_slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.5))
    ty_frame = ty_box.text_frame
    ty_frame.word_wrap = True
    ty_p = ty_frame.paragraphs[0]
    ty_p.alignment = PP_ALIGN.CENTER
    # Bilingual: if title contains Arabic characters use Arabic, else English
    is_arabic = any(ord(c) > 0x600 for c in title)
    ty_p.text = "شكراً" if is_arabic else "THANK YOU"
    ty_p.font.name = tf
    ty_p.font.size = Pt(54)
    ty_p.font.bold = True
    ty_p.font.color.rgb = tc

    # Subtitle line
    sub_box = thank_you_slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.333), Inches(1.2))
    sub_frame = sub_box.text_frame
    sub_frame.word_wrap = True
    sub_p = sub_frame.paragraphs[0]
    sub_p.alignment = PP_ALIGN.CENTER
    sub_p.text = title
    sub_p.font.name = bf
    sub_p.font.size = Pt(22)
    sub_p.font.color.rgb = bc

    # Divider line under title
    try:
        from pptx.util import Emu
        divider = thank_you_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.3), Inches(4.333), Pt(3)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = cc
        divider.line.fill.background()
    except Exception:
        pass

    # ── PAGE NUMBERS ─────────────────────────────────────────────────────────
    # One final pass over every slide already built, rather than threading
    # numbering through every individual slide-creation call above — safer
    # (touches nothing about how slides are built) and simpler (one place
    # to change later). Skips the very first slide (the cover/title slide),
    # which conventionally doesn't carry a page number.
    total_slides = len(list(prs.slides))
    for i, s in enumerate(prs.slides, start=1):
        if i == 1:
            continue  # skip the cover slide
        add_page_number(s, i, total_slides, RGBColor(150, 150, 150), "Calibri")

    prs.save(output_path)

    # Clean up temporary images AFTER prs.save()
    for img_p in temp_images:
        try:
            if os.path.exists(img_p):
                os.remove(img_p)
        except Exception:
            pass

    return output_path
