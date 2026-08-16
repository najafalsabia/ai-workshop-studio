import sys
import os
import json
import copy
import zipfile
import io
from pathlib import Path
import streamlit as st

# Setup system path to find backend files and load backend .env
root_dir = Path(__file__).resolve().parent
if root_dir.name == "Frontend":
    backend_dir = root_dir.parent / "Backend"
else:
    backend_dir = root_dir / "Backend"
sys.path.insert(0, str(backend_dir))

from dotenv import load_dotenv
load_dotenv(backend_dir / ".env")

# Verify keys are loaded
if not os.environ.get("GEMINI_API_KEY") or not os.environ.get("TAVILY_API_KEY"):
    st.warning("⚠️ Warning: GEMINI_API_KEY or TAVILY_API_KEY is not set. Please set them in your environment or Backend/.env file.")

# Import backend functions
try:
    from idea_agent import generate_titles
    from plan_builder import build_plan
    from content_generator import generate_content, generate_content_with_sources
    from activities_generator import extract_lab_contexts, suggest_lab_types, generate_one_lab
    from quiz_generator import generate_quiz
    from feedback_loop import build_revision_prompt, revise_single_slide
    from llm_client import ask_llm_for_json
    from notebook_builder import save_all_lab_notebooks
    from quiz_doc_builder import save_quiz_docx
    from kahoot_export import export_quiz_to_kahoot_xlsx
    from quality_checklist import run_quality_checklist
except ImportError as e:
    st.error(f"Failed to import backend modules. Ensure they are located in {backend_dir}. Error: {e}")
    st.stop()

# Translation Dictionary
T = {
    "English": {
        "step_1": "1. Parameter Form",
        "step_2": "2. Select Title",
        "step_3": "3. Workshop Plan",
        "step_4": "4. Slide Content",
        "step_5": "5. Lab Generator",
        "step_6": "6. Quiz",
        "step_7": "7. Quality Check & Export",
        
        "app_title": "AI Workshop Studio 🎨",
        "app_subtitle": "Design premium, hands-on technical workshops in minutes",
        
        "sidebar_config": "Sidebar Settings ⚙️",
        "lang_toggle": "UI Language / اللغة",
        "cycle_config": "Lab Cycle Settings 🧪",
        "use_cycle": "Use explain ➔ lab ➔ break cycle",
        "explain_mins": "Explain Duration (minutes)",
        "lab_mins": "Lab Duration (minutes)",
        "break_mins": "Break Duration (minutes)",
        
        "form_title": "Workshop Parameter Form 📝",
        "audience": "Target Audience",
        "age": "Age Group",
        "duration": "Workshop Duration (e.g., 3 hours, 90 minutes)",
        "goal": "Primary Learning Goal",
        "notes": "Extra Customization Notes",
        "idea_mode": "Idea Selection Mode",
        "have_idea": "I have an idea",
        "need_inspiration": "I need inspiration",
        "idea_input": "Your Workshop Idea",
        "idea_input_placeholder": "E.g., An introduction to prompt engineering for writers...",
        
        "btn_generate_titles": "Generate Title Suggestions ✨",
        "btn_next": "Next Step ➔",
        "btn_back": "⬅ Prev Step",
        "btn_start_over": "Start Over 🔄",
        
        "step2_title": "Select a Workshop Title 🏆",
        "step2_desc": "Choose one of the AI-suggested titles based on web research and trends:",
        "choose_btn": "Choose this Title",
        "refine_titles": "Refine Titles with AI",
        "refine_placeholder": "E.g., Make them more advanced, focus on marketing...",
        "btn_refine": "Apply Feedback 🔄",
        
        "step3_title": "Workshop Plan & Outline 📋",
        "step3_desc": "Review the generated learning objectives and timeline skeleton:",
        "objectives": "Learning Objectives",
        "outline": "Workshop Timeline",
        "refine_plan": "Refine Plan with AI",
        "refine_plan_placeholder": "E.g., Add a section on API keys, make the break longer...",
        "btn_generate_content": "Generate Slide Content 🚀",
        
        "step4_title": "Slide Content Editor 📝",
        "step4_desc": "Edit slide text blocks manually below, or ask AI to revise the whole deck:",
        "refine_content": "Refine Content with AI",
        "refine_content_placeholder": "E.g., Add more code examples, use simpler terms...",
        "btn_save_edits": "Save Edits & Continue ➔",
        
        "step5_title": "Hands-on Labs & Notebooks 🧪",
        "step5_desc": "Confirm lab types and generate custom lab materials for trainees:",
        "lab_slot": "Lab Slot",
        "covers": "Covers",
        "suggested_type": "AI Suggested Type",
        "confirm_type": "Confirm Type",
        "custom_notes": "Customization notes for this lab",
        "btn_generate_labs": "Generate Full Labs ⚙️",
        "lab_materials": "Generated Lab Materials",
        "trainee_preview": "Trainee Notebook Preview",
        "solution_preview": "Solution Notebook Preview",
        "instructor_notes": "Instructor Notes & Best Practices",
        "platforms": "Suggested Platforms",
        "refine_labs": "Refine Labs with AI",
        
        "step6_title": "Workshop Quiz & Export 🎓",
        "step6_desc": "Review the generated tiered quiz and export the workshop package:",
        "quiz_title": "Workshop Quiz",
        "export_btn": "Download Quiz JSON 📥",
        "refine_quiz": "Refine Quiz with AI",
        "correct_ans": "Correct Answer",
        
        "translate_btn": "Translate content to Arabic 🌐",
        "translate_success": "Content translated to Arabic successfully!",
        "translate_failed": "Failed to translate content.",
        "view_original": "View English Version 🇬🇧",
        "view_translated": "View Arabic Version 🇸🇦",
        
        "slide_options": "🎨 Slide Generation Options",
        "gen_mode_label": "Choose Slide Content Generation Mode:",
        "mode_ai_choice": "Generate from AI's Choice (Default)",
        "mode_mimic": "Mimic Uploaded Slide Reference",
        "upload_label": "Upload Slide Reference File (.json, .txt, or .pptx)",
    },
    "Arabic": {
        "step_1": "1. نموذج المعلومات",
        "step_2": "2. اختيار العنوان",
        "step_3": "3. خطة الورشة",
        "step_4": "4. محتوى الشرائح",
        "step_5": "5. المختبرات والأنشطة",
        "step_6": "6. الاختبار",
        "step_7": "7. فحص الجودة والتصدير",
        
        "app_title": "استوديو ورش عمل الذكاء الاصطناعي 🎨",
        "app_subtitle": "تصميم ورش عمل تقنية وتطبيقية مميزة خلال دقائق",
        
        "sidebar_config": "إعدادات الشريط الجانبي ⚙️",
        "lang_toggle": "لغة الواجهة / UI Language",
        "cycle_config": "إعدادات دورة المختبر 🧪",
        "use_cycle": "استخدام دورة (شرح ➔ مختبر ➔ استراحة)",
        "explain_mins": "مدة الشرح (بالدقائق)",
        "lab_mins": "مدة المختبر (بالدقائق)",
        "break_mins": "مدة الاستراحة (بالدقائق)",
        
        "form_title": "نموذج معلومات ورشة العمل 📝",
        "audience": "الجمهور المستهدف",
        "age": "الفئة العمرية",
        "duration": "مدة ورشة العمل (مثال: 3 ساعات، 90 دقيقة)",
        "goal": "هدف التعلم الأساسي",
        "notes": "ملاحظات تخصيص إضافية",
        "idea_mode": "وضع اختيار الفكرة",
        "have_idea": "لدي فكرة بالفعل",
        "need_inspiration": "أحتاج إلى إلهام",
        "idea_input": "فكرة ورشة العمل الخاصة بك",
        "idea_input_placeholder": "مثال: مقدمة في هندسة الأوامر للكتاب...",
        
        "btn_generate_titles": "توليد اقتراحات العناوين ✨",
        "btn_next": "الخطوة التالية ➔",
        "btn_back": "⬅ الخطوة السابقة",
        "btn_start_over": "البدء من جديد 🔄",
        
        "step2_title": "اختر عنوان ورشة العمل 🏆",
        "step2_desc": "اختر أحد العناوين المقترحة من الذكاء الاصطناعي بناءً على أبحاث الويب والاتجاهات الحديثة:",
        "choose_btn": "اختر هذا العنوان",
        "refine_titles": "تعديل العناوين بالذكاء الاصطناعي",
        "refine_placeholder": "مثال: اجعلها أكثر تقدمًا، ركز على التسويق...",
        "btn_refine": "تطبيق التعديلات 🔄",
        
        "step3_title": "خطة ورشة العمل وجدولها 📋",
        "step3_desc": "راجع الأهداف التعليمية والهيكل الزمني المقترح لورشة العمل:",
        "objectives": "الأهداف التعليمية",
        "outline": "الجدول الزمني لورشة العمل",
        "refine_plan": "تعديل الخطة بالذكاء الاصطناعي",
        "refine_plan_placeholder": "مثال: أضف قسمًا عن مفاتيح واجهة برمجة التطبيقات (API)، اجعل الاستراحة أطول...",
        "btn_generate_content": "توليد محتوى الشرائح 🚀",
        
        "step4_title": "محرر محتوى الشرائح 📝",
        "step4_desc": "قم بتعديل كتل نصوص الشرائح يدويًا أدناه، أو اطلب من الذكاء الاصطناعي مراجعة وتعديل كامل المحتوى:",
        "refine_content": "تعديل المحتوى بالذكاء الاصطناعي",
        "refine_content_placeholder": "مثال: أضف المزيد من أمثلة الكود، استخدم مصطلحات أبسط...",
        "btn_save_edits": "حفظ التعديلات والاستمرار ➔",
        
        "step5_title": "المختبرات التطبيقية ودفاتر الملاحظات 🧪",
        "step5_desc": "قم بتأكيد أنواع المختبرات وتوليد المواد التدريبية للمشاركين:",
        "lab_slot": "موقع المختبر",
        "covers": "يغطي مواضيع",
        "suggested_type": "النوع المقترح من الذكاء الاصطناعي",
        "confirm_type": "تأكيد النوع",
        "custom_notes": "ملاحظات تخصيص لهذا المختبر",
        "btn_generate_labs": "توليد المختبرات بالكامل ⚙️",
        "lab_materials": "المواد التدريبية المتولدة للمختبرات",
        "trainee_preview": "معاينة دفتر ملاحظات المتدرب",
        "solution_preview": "معاينة الحل الكامل",
        "instructor_notes": "ملاحظات المدرب وأفضل الممارسات",
        "platforms": "المنصات المقترحة",
        "refine_labs": "تعديل المختبرات بالذكاء الاصطناعي",
        
        "step6_title": "اختبار الورشة والتصدير 🎓",
        "step6_desc": "راجع الاختبار متعدد المستويات المتولد وقم بتصدير حزمة ورشة العمل:",
        "quiz_title": "اختبار ورشة العمل",
        "export_btn": "تحميل ملف الاختبار JSON 📥",
        "refine_quiz": "تعديل الاختبار بالذكاء الاصطناعي",
        "correct_ans": "الإجابة الصحيحة",
        
        "translate_btn": "ترجمة المحتوى إلى العربية 🌐",
        "translate_success": "تمت ترجمة المحتوى إلى العربية بنجاح!",
        "translate_failed": "فشلت ترجمة المحتوى.",
        "view_original": "عرض النسخة الإنجليزية 🇬🇧",
        "view_translated": "عرض النسخة العربية 🇸🇦",
        
        "slide_options": "🎨 نمط توليد الشرائح",
        "gen_mode_label": "اختر طريقة توليد محتوى الشرائح:",
        "mode_ai_choice": "التوليد التلقائي حسب اختيار الذكاء الاصطناعي (افتراضي)",
        "mode_mimic": "محاكاة مرجع الشرائح المرفوع",
        "upload_label": "ارفع ملف مرجع الشرائح (.json أو .txt أو .pptx)",
    }
}

MIMIC_CONTENT_PROMPT_TEMPLATE = """You are writing the full content for a technical workshop.

Workshop title: {title}

Learning objectives:
{objectives_list}

Outline (write content for EVERY section listed here, in the same order):
{outline_list}

This content will be shown on presentation slides, so brevity is a hard
requirement, not a style preference:
- Each block must be short enough to read at a glance. A "paragraph"
  block is 1-2 short sentences MAX — never a dense wall of text.
- A "bullet_list" item is a short phrase, not a full sentence.
- If a section has a lot to say, split it into MORE blocks (more slides),
  never into fewer, denser blocks. A speaker should be able to glance at
  a block and immediately know what to say next, not read paragraphs off
  the slide.

Use these block types as needed:
- "heading": a short section or subsection title
- "paragraph": 1-2 short sentences, explanatory
- "bullet_list": a list of short phrases (text is an array of strings)
- "code": a short code snippet (text is the code itself, plain text) —
  only when the section is actually about writing or reading code
- "image_placeholder": text describes what image should go here (no real image)

Write in English. Keep content specific and practical, matching the
section's description and duration — don't pad a 1-minute section with
20 minutes of reading material.

CRITICAL USER REQUIREMENT - MIMIC STRUCTURE & STYLE:
The user has provided an example of slides. You must MIMIC the block sequence, formatting style, and layout structure of the example slides below. Keep the same structure (e.g. if the example has heading, bullet_list, then code block, follow that sequence in your sections) but generate brand new technical content adapted to this workshop's goals and sections.

Here are the example slides to mimic:
{mimic_example}

CRITICAL: the "section" field in your reply must be copied EXACTLY,
character-for-character, from the outline above. Do NOT add the duration,
parentheses, or anything else to it — copy just the section name as-is.

Reply with ONLY this JSON, nothing else:
{{
  "content": [
    {{
      "section": "copied exactly from the outline, no additions",
      "blocks": [
        {{"block_id": "s1-b1", "type": "heading", "text": "..."}},
        {{"block_id": "s1-b2", "type": "paragraph", "text": "..."}}
      ]
    }}
  ]
}}
"""

# Page configuration
st.set_page_config(
    page_title="AI Workshop Studio",
    page_icon="🎨",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Disk State Persistence for Page Refresh Recovery
STATE_FILE = "temp_session_state.json"

def save_state_to_disk():
    state_keys = [
        "step", "language", "audience", "age", "duration", "goal", "notes",
        "idea_option", "idea_input", "slide_gen_mode_choice", "use_lab_cycle",
        "explain_minutes", "lab_minutes", "break_minutes",
        "titles_result", "chosen_title", "plan_result", "content_result",
        "lab_contexts", "lab_suggestions", "confirmed_lab_types", "lab_customizations",
        "labs_result", "quiz_result", "quiz_approved", "quiz_approved_snapshot",
        "quality_result", "quality_feedback_texts",
        "titles_result_ar", "plan_result_ar", "content_result_ar", "labs_result_ar", "quiz_result_ar",
        "view_language", "num_quiz_questions", "slide_style", "style_application"
    ]
    data = {}
    for key in state_keys:
        if key in st.session_state:
            data[key] = st.session_state[key]
    # Dynamically save per-section style keys
    for key in list(st.session_state.keys()):
        if key.startswith("style_section_"):
            data[key] = st.session_state[key]
    try:
        with open(STATE_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception:
        pass

def load_state_from_disk():
    if os.path.exists(STATE_FILE):
        try:
            with open(STATE_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            for key, val in data.items():
                st.session_state[key] = val
        except Exception:
            pass

def clear_state_disk():
    if os.path.exists(STATE_FILE):
        try:
            os.remove(STATE_FILE)
        except Exception:
            pass

# Monkeypatch st.rerun to ensure state is saved before immediate script aborts
_original_rerun = st.rerun
def custom_rerun():
    save_state_to_disk()
    _original_rerun()
st.rerun = custom_rerun

# Initialize Session State
def init_state(force=False):
    if force:
        clear_state_disk()
        for key in list(st.session_state.keys()):
            try:
                del st.session_state[key]
            except Exception:
                pass
    else:
        # Only load from disk on a fresh browser page load (when st.session_state is empty)
        # to prevent overwriting active user input changes during session reruns.
        if len(st.session_state.keys()) == 0:
            load_state_from_disk()
        
    default_states = {
        "step": 1,
        "language": "English",
        
        # Form Parameters
        "audience": "university computer science students",
        "age": "18-24",
        "duration": "3 hours",
        "goal": "teach practical use of AI coding assistants",
        "notes": "should feel hands-on, not just slides",
        "idea_option": "I need inspiration",
        "idea_input": "",
        "slide_gen_mode_choice": "ai",
        
        # Lab Cycle Settings
        "use_lab_cycle": True,
        "explain_minutes": 30,
        "lab_minutes": 20,
        "break_minutes": 10,
        
        # Data Persistence (English)
        "titles_result": None,
        "chosen_title": None,
        "plan_result": None,
        "content_result": None,
        "lab_contexts": None,
        "lab_suggestions": None,
        "confirmed_lab_types": {},
        "lab_customizations": {},
        "labs_result": None,
        "quiz_result": None,
        # True only after the trainer explicitly clicks "Approve Final Quiz".
        # quiz_approved_snapshot is a FROZEN copy of the quiz at the moment
        # of approval — every export (docx, Kahoot, JSON) reads from this
        # snapshot, never from the live quiz_result, so an export can never
        # accidentally contain a stale or mid-edit version.
        "quiz_approved": False,
        "quiz_approved_snapshot": None,

        # Step 7: final quality checklist result, and per-category feedback
        # text the trainer can review/edit before it's sent as a revision
        # request — Step 7 drafts the starting text, but nothing is ever
        # sent without the trainer seeing and approving it first (see the
        # Step 7 section below).
        "quality_result": None,
        "quality_feedback_texts": {},
        
        # Data Persistence (Arabic Translations)
        "titles_result_ar": None,
        "plan_result_ar": None,
        "content_result_ar": None,
        "labs_result_ar": None,
        "quiz_result_ar": None,
        
        # Language preference for generated content view (when UI is Arabic)
        "view_language": "English"
    }
    for key, val in default_states.items():
        if force or key not in st.session_state or (isinstance(st.session_state[key], str) and st.session_state[key].strip() == ""):
            try:
                st.session_state[key] = val
            except Exception:
                # If key is bound to a widget, delete it so it defaults to the initial value parameter
                try:
                    del st.session_state[key]
                except Exception:
                    pass

init_state()

# Inject Premium Custom Styles
st.markdown("""
<style>
    /* Dark glassmorphic design system */
    .title-card {
        background-color: #161b22;
        border: 1px solid #30363d;
        border-radius: 12px;
        padding: 24px;
        margin-bottom: 20px;
        box-shadow: 0 4px 15px rgba(0, 0, 0, 0.25);
        transition: all 0.3s cubic-bezier(0.16, 1, 0.3, 1);
    }
    .title-card:hover {
        border-color: #7928ca;
        transform: translateY(-2px);
        box-shadow: 0 8px 24px rgba(121, 40, 202, 0.25);
    }
    .title-card, .title-card *, .timeline-card, .timeline-card * {
        color: #f0f6fc !important;
    }
    .title-card h4 {
        color: #58a6ff !important;
        margin-top: 0 !important;
        margin-bottom: 8px !important;
    }
    .title-card p {
        color: #c9d1d9 !important;
    }
    
    /* Horizontal progress step bar */
    .step-bar {
        display: flex;
        justify-content: space-between;
        margin-bottom: 35px;
        gap: 8px;
    }
    .step-unit {
        flex: 1;
        text-align: center;
        padding: 12px 6px;
        font-size: 14px;
        font-weight: 600;
        border-bottom: 4px solid #30363d;
        color: #8b949e;
        transition: all 0.2s ease;
    }
    .step-unit.active {
        border-bottom: 4px solid #7928ca;
        color: #f0f6fc;
        background: rgba(121, 40, 202, 0.05);
    }
    .step-unit.completed {
        border-bottom: 4px solid #238636;
        color: #3fb950;
    }
    
    /* Visual timeline design */
    .timeline-container {
        border-left: 3px solid #30363d;
        padding-left: 24px;
        margin-left: 14px;
        position: relative;
    }
    .timeline-item {
        margin-bottom: 30px;
        position: relative;
    }
    .timeline-dot {
        position: absolute;
        left: -33px;
        top: 4px;
        background-color: #7928ca;
        width: 16px;
        height: 16px;
        border-radius: 50%;
        border: 3px solid #0d1117;
        box-shadow: 0 0 8px rgba(121, 40, 202, 0.8);
    }
    .timeline-card {
        background-color: #161b22;
        border: 1px solid #30363d;
        border-radius: 10px;
        padding: 16px;
        box-shadow: 0 3px 10px rgba(0,0,0,0.15);
    }
    
    /* Micro animations and custom fonts */
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');
    body, p, h1, h2, h3, h4, h5, h6 {
        font-family: 'Inter', -apple-system, sans-serif;
    }
</style>
""", unsafe_allow_html=True)

# RTL Support CSS Injector
if st.session_state.language == "Arabic":
    st.markdown("""
    <style>
        .main .block-container {
            direction: rtl !important;
            text-align: right !important;
        }
        .stTextInput input, .stTextArea textarea, .stSelectbox select {
            direction: rtl !important;
            text-align: right !important;
        }
        .timeline-container {
            border-left: none !important;
            border-right: 3px solid #30363d !important;
            padding-left: 0 !important;
            padding-right: 24px !important;
            margin-left: 0 !important;
            margin-right: 14px !important;
        }
        .timeline-dot {
            left: auto !important;
            right: -33px !important;
        }
    </style>
    """, unsafe_allow_html=True)

# Helper: Translate JSON content using Gemini API
def translate_data_to_arabic(data: dict, context_type: str) -> dict:
    prompt = f"""You are a professional technical translator translating workshop materials into Arabic.
Translate the following JSON content from English to Arabic.
Keep the JSON keys, structure, types, and values exactly the same. Only translate the text values.
For code blocks or code cells in notebooks, do NOT translate the actual code itself, only translate comments and explanations.
Ensure technical accuracy and natural phrasing.

Context: {context_type}

JSON Content:
{json.dumps(data, indent=2, ensure_ascii=False)}

Reply with ONLY the translated JSON, nothing else."""
    
    return ask_llm_for_json(prompt)

# Helper: Clear dynamic block keys
def clear_block_keys():
    keys_to_del = [k for k in st.session_state.keys() if k.startswith("block_")]
    for k in keys_to_del:
        del st.session_state[k]

# Helper: Run dynamic revision loop
def revise_output(current_output: dict, feedback: str, context_description: str) -> dict:
    prompt = build_revision_prompt(current_output, feedback, context_description)
    return ask_llm_for_json(prompt)

# Helper: Extract slides text from PPTX file
def extract_pptx_slides_text(pptx_file_stream) -> str:
    from pptx import Presentation
    prs = Presentation(pptx_file_stream)
    slides_content = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        slide_text.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        if paragraph.level > 0:
                            slide_text.append(f"  - {para_text}")
                        else:
                            slide_text.append(f"  {para_text}")
        slides_content.append("\n".join(slide_text))
        
    return "\n\n".join(slides_content)

# Sidebar Configuration
lang = st.session_state.language
with st.sidebar:
    st.header(T[lang]["sidebar_config"])
    
    # UI Language Switch
    st.selectbox(
        T[lang]["lang_toggle"],
        options=["English", "Arabic"],
        key="language"
    )
    
    # Force redraw immediately on language change
    if st.session_state.language != lang:
        st.session_state.view_language = st.session_state.language
        for k in ["view_language_sel_2", "view_language_sel_3", "view_language_sel_4", "view_language_sel_5", "view_language_sel_6"]:
            if k in st.session_state:
                del st.session_state[k]
        st.rerun()
    
    st.markdown("---")
    
    # Lab Cycle settings
    st.subheader(T[lang]["cycle_config"])
    st.checkbox(T[lang]["use_cycle"], key="use_lab_cycle")
    
    if st.session_state.use_lab_cycle:
        st.number_input(T[lang]["explain_mins"], min_value=1, max_value=240, key="explain_minutes")
        st.number_input(T[lang]["lab_mins"], min_value=1, max_value=240, key="lab_minutes")
        st.number_input(T[lang]["break_mins"], min_value=1, max_value=240, key="break_minutes")
        
    if "slide_style" not in st.session_state:
        st.session_state.slide_style = "Clean & Minimal"
    if "style_application" not in st.session_state:
        st.session_state.style_application = "All slides"
            
    st.markdown("---")

    st.markdown("### 📂 " + ("My Workshops" if lang == "English" else "ورشي المحفوظة"))
    try:
        from workshop_db import list_workshops, load_workshop, delete_workshop
        saved_workshops = list_workshops()
    except Exception as e:
        saved_workshops = []
        st.caption(f"Could not load saved workshops: {e}")

    if not saved_workshops:
        st.caption("No saved workshops yet." if lang == "English" else "ما فيه ورش محفوظة بعد.")
    else:
        for w in saved_workshops:
            date_only = w["created_at"][:10]
            with st.expander(f"{w['title']} — {date_only}"):
                st.caption(f"{w['audience']} · {w['duration']}")
                col_load, col_del = st.columns(2)
                with col_load:
                    if st.button(
                        "📂 Load" if lang == "English" else "📂 فتح",
                        key=f"load_workshop_{w['id']}",
                        use_container_width=True,
                    ):
                        full = load_workshop(w["id"])
                        st.session_state.chosen_title = full["title"]
                        st.session_state.audience = full["audience"]
                        st.session_state.age = full["age"]
                        st.session_state.duration = full["duration"]
                        st.session_state.plan_result = full["plan"]
                        st.session_state.content_result = full["content"]
                        st.session_state.labs_result = full["labs"]
                        st.session_state.quiz_result = full["quiz"]
                        st.session_state.plan_result_ar = None
                        st.session_state.content_result_ar = None
                        st.session_state.labs_result_ar = None
                        st.session_state.quiz_result_ar = None
                        st.session_state.quiz_approved = True
                        st.session_state.quiz_approved_snapshot = full["quiz"]
                        st.session_state.quality_result = None
                        st.session_state.step = 7
                        st.rerun()
                with col_del:
                    if st.button("🗑️", key=f"delete_workshop_{w['id']}", use_container_width=True):
                        delete_workshop(w["id"])
                        st.rerun()

    st.markdown("---")
    
    # Reset application
    if st.button(T[lang]["btn_start_over"], type="secondary", use_container_width=True):
        init_state(force=True)
        clear_block_keys()
        st.rerun()

# Renders the beautiful progress step bar at the top
def render_progress_bar():
    steps = [
        T[lang]["step_1"],
        T[lang]["step_2"],
        T[lang]["step_3"],
        T[lang]["step_4"],
        T[lang]["step_5"],
        T[lang]["step_6"],
        T[lang]["step_7"],
    ]
    
    st.markdown("<div class='step-bar'>", unsafe_allow_html=True)
    cols = st.columns(len(steps))
    current_step = st.session_state.step
    for idx, name in enumerate(steps, 1):
        with cols[idx - 1]:
            if idx == current_step:
                st.markdown(f"<div class='step-unit active'>{name}</div>", unsafe_allow_html=True)
            elif idx < current_step:
                st.markdown(f"<div class='step-unit completed'>✓ {name}</div>", unsafe_allow_html=True)
            else:
                st.markdown(f"<div class='step-unit'>{name}</div>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

# Main Application Window
st.title(T[lang]["app_title"])
st.caption(T[lang]["app_subtitle"])
render_progress_bar()

# STEP 1: Parameter Form Input
if st.session_state.step == 1:
    st.header(T[lang]["form_title"])
    
    # Mode toggle at top
    st.radio(T[lang]["idea_mode"], options=[T[lang]["need_inspiration"], T[lang]["have_idea"]], key="idea_option", horizontal=True)
    
    st.markdown("---")
    
    is_inspiration = st.session_state.idea_option == T[lang]["need_inspiration"]
    
    if is_inspiration:
        # ── INSPIRATION MODE ──
        # Field selector – technology domains only, plus free-text "Other"
        OTHER_LABEL_EN = "Other (describe your field below)"
        OTHER_LABEL_AR = "أخرى (اكتب مجالك أدناه)"
        fields_en = [
            "Technology & Software",
            "Data Science & AI",
            "Cybersecurity",
            "Cloud & DevOps",
            "Engineering & Hardware",
            OTHER_LABEL_EN,
        ]
        fields_ar = [
            "التكنولوجيا والبرمجيات",
            "علوم البيانات والذكاء الاصطناعي",
            "الأمن السيبراني",
            "الحوسبة السحابية و DevOps",
            "الهندسة والأجهزة",
            OTHER_LABEL_AR,
        ]
        fields = fields_en if lang == "English" else fields_ar
        other_label = OTHER_LABEL_EN if lang == "English" else OTHER_LABEL_AR

        st.selectbox(
            "📂 " + ("Select a Field" if lang == "English" else "اختر مجالاً"),
            options=fields,
            key="inspiration_field"
        )

        # Show free-text input when "Other" is chosen
        is_other_field = st.session_state.get("inspiration_field", "") == other_label
        if is_other_field:
            st.text_input(
                "✏️ " + ("Describe your field or topic" if lang == "English" else "صف مجالك أو موضوعك"),
                placeholder=("e.g. Quantum Computing, Blockchain, AR/VR, Bioinformatics…" if lang == "English" else "مثال: الحوسبة الكمية، البلوك تشين، الواقع المعزز…"),
                key="inspiration_field_other"
            )
        
        # Shared inputs
        col1, col2, col3 = st.columns(3)
        with col1:
            st.text_input(T[lang]["audience"], value=st.session_state.get("audience", ""), key="audience")
        with col2:
            st.text_input(T[lang]["age"], value=st.session_state.get("age", ""), key="age")
        with col3:
            st.text_input(T[lang]["duration"], value=st.session_state.get("duration", ""), key="duration")
        
        st.text_area(T[lang]["notes"], value=st.session_state.get("notes", ""), key="notes", height=80)
        
        # Topic inspiration cards filtered by field
        all_topics = {
            "Technology & Software": [
                {"emoji": "🤖", "title": "AI-Powered Coding Assistants", "goal": "Learn to use AI tools like GitHub Copilot and ChatGPT for faster, smarter coding"},
                {"emoji": "📱", "title": "Mobile App Development with Flutter", "goal": "Build beautiful cross-platform mobile apps from scratch using Dart and Flutter"},
                {"emoji": "🌐", "title": "Full-Stack Web Development", "goal": "Build modern web apps with React, Node.js, and REST APIs end-to-end"},
            ],
            "Data Science & AI": [
                {"emoji": "🧠", "title": "Machine Learning from Scratch", "goal": "Build and train ML models step-by-step without black-box libraries"},
                {"emoji": "📊", "title": "Data Visualization with Python", "goal": "Build interactive dashboards using Plotly, Streamlit, and real datasets"},
                {"emoji": "🎨", "title": "Prompt Engineering Masterclass", "goal": "Master the art of writing effective prompts for large language models"},
            ],
            "Business & Entrepreneurship": [
                {"emoji": "🚀", "title": "Startup Pitch & Business Model", "goal": "Craft a compelling startup pitch deck and validate your business model canvas"},
                {"emoji": "💼", "title": "Product Management Essentials", "goal": "Learn to prioritize features, write user stories, and manage product roadmaps"},
                {"emoji": "📋", "title": "Lean Startup Methodology", "goal": "Apply build-measure-learn cycles to validate ideas with minimal resources"},
            ],
            "Design & Creative": [
                {"emoji": "🎨", "title": "UX/UI Design Thinking Workshop", "goal": "Apply design thinking methodology to create user-centered digital products"},
                {"emoji": "✏️", "title": "Figma for Beginners", "goal": "Design professional UI mockups and interactive prototypes using Figma"},
                {"emoji": "🖌️", "title": "Brand Identity & Visual Design", "goal": "Create cohesive brand guidelines including logos, typography, and color systems"},
            ],
            "Cybersecurity": [
                {"emoji": "🔒", "title": "Cybersecurity Essentials", "goal": "Understand threat landscapes, penetration testing, and secure coding practices"},
                {"emoji": "🛡️", "title": "Ethical Hacking & Bug Bounty", "goal": "Learn reconnaissance, exploitation, and responsible disclosure workflows"},
                {"emoji": "🔐", "title": "Zero Trust Security Architecture", "goal": "Design and implement zero trust network security models for modern organizations"},
            ],
            "Cloud & DevOps": [
                {"emoji": "☁️", "title": "Cloud Architecture Fundamentals", "goal": "Design scalable cloud infrastructure using AWS, Azure, or GCP services"},
                {"emoji": "🐳", "title": "Docker & Kubernetes Workshop", "goal": "Containerize applications and orchestrate deployments with Kubernetes"},
                {"emoji": "⚙️", "title": "CI/CD Pipeline Automation", "goal": "Build automated testing and deployment pipelines with GitHub Actions and Jenkins"},
            ],
            "Marketing & Growth": [
                {"emoji": "📈", "title": "Digital Marketing Strategy", "goal": "Create data-driven marketing campaigns across SEO, social media, and email"},
                {"emoji": "📣", "title": "Content Marketing & Storytelling", "goal": "Craft compelling content strategies that drive engagement and conversions"},
                {"emoji": "📊", "title": "Growth Hacking Techniques", "goal": "Apply rapid experimentation frameworks to accelerate user acquisition"},
            ],
            "Leadership & Soft Skills": [
                {"emoji": "🎤", "title": "Public Speaking & Presentation Skills", "goal": "Master storytelling, body language, and slide design for impactful presentations"},
                {"emoji": "🤝", "title": "Agile & Scrum Project Management", "goal": "Run effective sprints, standups, and retrospectives for software teams"},
                {"emoji": "🧭", "title": "Leadership & Team Building", "goal": "Develop leadership styles, conflict resolution, and team motivation strategies"},
            ],
            "Engineering & Hardware": [
                {"emoji": "🔧", "title": "Robotics & IoT with Arduino", "goal": "Build and program sensor-driven robots using Arduino boards and C++"},
                {"emoji": "⚡", "title": "Embedded Systems Programming", "goal": "Program microcontrollers and design real-time embedded applications"},
                {"emoji": "🏗️", "title": "3D Printing & CAD Design", "goal": "Design 3D models in CAD software and bring them to life with 3D printing"},
            ],
            "Education & Training": [
                {"emoji": "📚", "title": "Instructional Design Workshop", "goal": "Design effective training programs using ADDIE and SAM frameworks"},
                {"emoji": "🎮", "title": "Gamification in Learning", "goal": "Apply game mechanics to training and education for higher engagement"},
                {"emoji": "🧑‍🏫", "title": "Train the Trainer Program", "goal": "Equip trainers with facilitation skills, assessment design, and delivery techniques"},
            ],
        }
        # Arabic topic mapping
        all_topics_ar = {
            "التكنولوجيا والبرمجيات": [
                {"emoji": "🤖", "title": "مساعدات الترميز بالذكاء الاصطناعي", "goal": "تعلم استخدام أدوات الذكاء الاصطناعي للترميز الأسرع"},
                {"emoji": "📱", "title": "تطوير تطبيقات الجوال بـ Flutter", "goal": "بناء تطبيقات جوال جميلة ومتعددة المنصات من الصفر"},
                {"emoji": "🌐", "title": "تطوير الويب الشامل", "goal": "بناء تطبيقات ويب حديثة باستخدام React و Node.js"},
            ],
            "علوم البيانات والذكاء الاصطناعي": [
                {"emoji": "🧠", "title": "تعلم الآلة من الصفر", "goal": "بناء وتدريب نماذج تعلم الآلة خطوة بخطوة"},
                {"emoji": "📊", "title": "تصور البيانات باستخدام بايثون", "goal": "بناء لوحات معلومات تفاعلية باستخدام بيانات حقيقية"},
                {"emoji": "🎨", "title": "هندسة الأوامر النصية", "goal": "إتقان كتابة أوامر فعالة لنماذج اللغة الكبيرة"},
            ],
            "الأعمال وريادة الأعمال": [
                {"emoji": "🚀", "title": "العرض التقديمي ونموذج العمل", "goal": "صياغة عرض تقديمي مقنع والتحقق من صحة نموذج العمل"},
                {"emoji": "💼", "title": "أساسيات إدارة المنتجات", "goal": "تعلم ترتيب الأولويات وكتابة قصص المستخدم وإدارة خارطة الطريق"},
                {"emoji": "📋", "title": "منهجية الشركة الناشئة الرشيقة", "goal": "تطبيق دورات البناء-القياس-التعلم للتحقق من الأفكار"},
            ],
            "التصميم والإبداع": [
                {"emoji": "🎨", "title": "ورشة التفكير التصميمي", "goal": "تطبيق منهجية التفكير التصميمي لإنشاء منتجات رقمية"},
                {"emoji": "✏️", "title": "فيجما للمبتدئين", "goal": "تصميم نماذج واجهات احترافية وتفاعلية باستخدام Figma"},
                {"emoji": "🖌️", "title": "هوية العلامة التجارية والتصميم البصري", "goal": "إنشاء إرشادات علامة تجارية متماسكة"},
            ],
            "الأمن السيبراني": [
                {"emoji": "🔒", "title": "أساسيات الأمن السيبراني", "goal": "فهم مشهد التهديدات واختبار الاختراق"},
                {"emoji": "🛡️", "title": "القرصنة الأخلاقية ومكافآت الثغرات", "goal": "تعلم الاستطلاع والاستغلال والإفصاح المسؤول"},
                {"emoji": "🔐", "title": "بنية أمان الثقة الصفرية", "goal": "تصميم وتنفيذ نماذج أمان الثقة الصفرية"},
            ],
            "الحوسبة السحابية و DevOps": [
                {"emoji": "☁️", "title": "أساسيات الحوسبة السحابية", "goal": "تصميم بنية سحابية قابلة للتوسع"},
                {"emoji": "🐳", "title": "ورشة Docker و Kubernetes", "goal": "حاويات التطبيقات وتنظيم النشر"},
                {"emoji": "⚙️", "title": "أتمتة خطوط CI/CD", "goal": "بناء خطوط اختبار ونشر آلية"},
            ],
            "التسويق والنمو": [
                {"emoji": "📈", "title": "استراتيجية التسويق الرقمي", "goal": "إنشاء حملات تسويقية مبنية على البيانات"},
                {"emoji": "📣", "title": "تسويق المحتوى والسرد القصصي", "goal": "صياغة استراتيجيات محتوى مقنعة"},
                {"emoji": "📊", "title": "تقنيات قرصنة النمو", "goal": "تطبيق أطر التجريب السريع لتسريع اكتساب المستخدمين"},
            ],
            "القيادة والمهارات الشخصية": [
                {"emoji": "🎤", "title": "مهارات العرض والتحدث أمام الجمهور", "goal": "إتقان فن السرد القصصي ولغة الجسد"},
                {"emoji": "🤝", "title": "إدارة المشاريع بأسلوب أجايل", "goal": "إدارة سباقات العمل والاجتماعات اليومية بفعالية"},
                {"emoji": "🧭", "title": "القيادة وبناء الفريق", "goal": "تطوير أساليب القيادة وحل النزاعات وتحفيز الفريق"},
            ],
            "الهندسة والأجهزة": [
                {"emoji": "🔧", "title": "الروبوتات وإنترنت الأشياء مع أردوينو", "goal": "بناء وبرمجة روبوتات تعمل بالمستشعرات"},
                {"emoji": "⚡", "title": "برمجة الأنظمة المدمجة", "goal": "برمجة المتحكمات الدقيقة وتصميم تطبيقات مدمجة"},
                {"emoji": "🏗️", "title": "الطباعة ثلاثية الأبعاد وتصميم CAD", "goal": "تصميم نماذج ثلاثية الأبعاد وطباعتها"},
            ],
            "التعليم والتدريب": [
                {"emoji": "📚", "title": "ورشة التصميم التعليمي", "goal": "تصميم برامج تدريبية فعالة"},
                {"emoji": "🎮", "title": "التلعيب في التعلم", "goal": "تطبيق آليات الألعاب في التدريب والتعليم"},
                {"emoji": "🧑‍🏫", "title": "برنامج تدريب المدربين", "goal": "تزويد المدربين بمهارات التيسير والتقييم"},
            ],
        }
        
        topics_map = all_topics if lang == "English" else all_topics_ar
        selected_field = st.session_state.get("inspiration_field", fields[0])
        field_topics = topics_map.get(selected_field, [])

        # Only show inspiration cards when the field has presets (not for "Other")
        if field_topics:
            st.markdown(f"<p style='font-size:13px; color:#888; margin-top:10px; margin-bottom:6px; font-family:sans-serif;'>{'💡 Click a topic to auto-fill, or just hit Generate:' if lang == 'English' else '💡 اضغط على موضوع لملء النموذج، أو اضغط توليد مباشرة:'}</p>", unsafe_allow_html=True)
            topic_cols = st.columns(len(field_topics))
            for t_idx, topic in enumerate(field_topics):
                with topic_cols[t_idx]:
                    if st.button(f"{topic['emoji']} {topic['title']}", key=f"insp_{t_idx}", use_container_width=True):
                        st.session_state.goal = topic["goal"]
                        st.session_state.idea_input = topic["title"]
                        st.rerun()
        
        # Show the auto-filled goal (editable)
        st.text_area(T[lang]["goal"], value=st.session_state.get("goal", ""), key="goal", height=80)
        
    else:
        # ── I HAVE AN IDEA MODE ──
        st.text_input(T[lang]["idea_input"], placeholder=T[lang]["idea_input_placeholder"], key="idea_input")
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.text_input(T[lang]["audience"], value=st.session_state.get("audience", ""), key="audience")
        with col2:
            st.text_input(T[lang]["age"], value=st.session_state.get("age", ""), key="age")
        with col3:
            st.text_input(T[lang]["duration"], value=st.session_state.get("duration", ""), key="duration")
        
        st.text_area(T[lang]["goal"], value=st.session_state.get("goal", ""), key="goal", height=80)
        st.text_area(T[lang]["notes"], value=st.session_state.get("notes", ""), key="notes", height=80)
    
    # Generate Title Suggestions button (shared by both modes)
    if st.button(T[lang]["btn_generate_titles"], type="primary", use_container_width=True):
        from plan_builder import parse_duration_to_minutes
        
        parsed_mins = parse_duration_to_minutes(st.session_state.duration)
        if parsed_mins is None:
            st.error("❌ " + ("Could not understand duration. Try formats like '3 hours' or '120 minutes'." if lang == "English" else "لم نتمكن من فهم المدة المكتوبة. يرجى كتابتها بصيغة مثل '3 ساعات' أو '120 دقيقة'."))
        else:
            notes_payload = st.session_state.get("notes", "")
            if is_inspiration:
                selected_field_val = st.session_state.get("inspiration_field", "")
                # If the user chose "Other", use their custom text instead of the generic label
                _other_lbl = OTHER_LABEL_EN if lang == "English" else OTHER_LABEL_AR
                if selected_field_val == _other_lbl:
                    selected_field_val = st.session_state.get("inspiration_field_other", "").strip() or selected_field_val
                notes_payload = f"Workshop Field/Domain: {selected_field_val}\n\n{notes_payload}"
            if st.session_state.idea_option == T[lang]["have_idea"] and st.session_state.get("idea_input"):
                notes_payload = f"User Specific Idea: {st.session_state.idea_input}\n\n{notes_payload}"
            elif is_inspiration and st.session_state.get("idea_input"):
                notes_payload = f"Inspiration Topic: {st.session_state.idea_input}\n\n{notes_payload}"
            
            with st.spinner("Searching latest trends and generating premium titles..." if lang == "English" else "جاري البحث في أحدث الاتجاهات وتوليد عناوين مميزة..."):
                try:
                    res = generate_titles(
                        audience=st.session_state.audience,
                        age=st.session_state.age,
                        duration=st.session_state.duration,
                        goal=st.session_state.get("goal", ""),
                        notes=notes_payload
                    )
                    st.session_state.titles_result = res
                    st.session_state.chosen_title = None
                    st.session_state.plan_result = None
                    st.session_state.content_result = None
                    st.session_state.labs_result = None
                    st.session_state.quiz_result = None
                    st.session_state.titles_result_ar = None
                    st.session_state.plan_result_ar = None
                    st.session_state.content_result_ar = None
                    st.session_state.labs_result_ar = None
                    st.session_state.quiz_result_ar = None
                    
                    st.session_state.step = 2
                    st.rerun()
                except Exception as e:
                    st.error(f"Error during title generation: {e}")
# STEP 2: Title Selection
elif st.session_state.step == 2:
    st.header(T[lang]["step2_title"])
    st.write(T[lang]["step2_desc"])
    
    # Handle Translation Option
    if lang == "Arabic":
        col_lang, _ = st.columns([1, 4])
        with col_lang:
            default_idx = 1 if st.session_state.view_language == "Arabic" else 0
            st.selectbox("View Language / عرض اللغة", options=[T[lang]["view_original"], T[lang]["view_translated"]], index=default_idx, key="view_language_sel_2")
            st.session_state.view_language = "Arabic" if st.session_state.view_language_sel_2 == T[lang]["view_translated"] else "English"
        
        # Dynamic translation trigger
        if st.session_state.view_language == "Arabic" and st.session_state.titles_result_ar is None:
            with st.spinner(T[lang]["translate_btn"]):
                try:
                    st.session_state.titles_result_ar = translate_data_to_arabic(st.session_state.titles_result, "Workshop titles list")
                except Exception as e:
                    st.error(f"{T[lang]['translate_failed']}: {e}")
                    st.session_state.view_language = "English"
    else:
        st.session_state.view_language = "English"

    # Fetch active titles
    active_titles_dict = st.session_state.titles_result_ar if (st.session_state.view_language == "Arabic" and st.session_state.titles_result_ar) else st.session_state.titles_result
    
    if active_titles_dict and "titles" in active_titles_dict:
        titles_list = active_titles_dict["titles"]
        for idx, t in enumerate(titles_list):
            st.markdown(f"""
            <div class="title-card">
                <h4>{t.get('title')}</h4>
                <p style="margin-bottom:0; font-size: 14px; opacity:0.85;">{t.get('why')}</p>
            </div>
            """, unsafe_allow_html=True)
            
            # Select button
            if st.button(f"{T[lang]['choose_btn']} #{idx+1}", key=f"select_title_{idx}", type="primary"):
                # Save the select English title as target reference
                eng_title = st.session_state.titles_result["titles"][idx]["title"]
                st.session_state.chosen_title = eng_title
                
                # Build the Plan
                notes_payload = st.session_state.notes
                if st.session_state.idea_option == T[lang]["have_idea"] and st.session_state.idea_input:
                    notes_payload = f"User Specific Idea: {st.session_state.idea_input}\n\n{notes_payload}"
                
                with st.spinner("Building custom plan outline..." if lang == "English" else "جاري بناء الهيكل الزمني للخطة..."):
                    try:
                        plan = build_plan(
                            title=eng_title,
                            audience=st.session_state.audience,
                            age=st.session_state.age,
                            duration=st.session_state.duration,
                            goal=st.session_state.goal,
                            notes=notes_payload,
                            use_lab_cycle=st.session_state.use_lab_cycle,
                            explain_minutes=st.session_state.explain_minutes,
                            lab_minutes=st.session_state.lab_minutes,
                            break_minutes=st.session_state.break_minutes
                        )
                        st.session_state.plan_result = plan
                        st.session_state.step = 3
                        st.rerun()
                    except Exception as e:
                        st.error(f"Error during plan generation: {e}")
                        
        st.markdown("---")
        
        # Interactive AI feedback revision
        st.subheader(T[lang]["refine_titles"])
        feedback = st.text_area(T[lang]["refine_titles"], placeholder=T[lang]["refine_placeholder"], label_visibility="collapsed")
        
        if st.button(T[lang]["btn_refine"]):
            if feedback.strip():
                with st.spinner("Revising titles..." if lang == "English" else "جاري مراجعة العناوين..."):
                    try:
                        notes_payload = st.session_state.notes
                        if st.session_state.idea_option == T[lang]["have_idea"] and st.session_state.idea_input:
                            notes_payload = f"User Specific Idea: {st.session_state.idea_input}\n\n{notes_payload}"
                        context = f"Title suggestions for audience={st.session_state.audience}, duration={st.session_state.duration}, goal={st.session_state.goal}, notes={notes_payload}"
                        
                        # Revise the English source title dict
                        revised = revise_output(st.session_state.titles_result, feedback, context)
                        st.session_state.titles_result = revised
                        st.session_state.titles_result_ar = None # invalidate cache
                        st.success("Titles revised successfully!" if lang == "English" else "تم مراجعة العناوين بنجاح!")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Failed to revise titles: {e}")
    else:
        st.error("No title options available. Go back and generate them again.")
        
    st.markdown("---")
    if st.button(T[lang]["btn_back"]):
        st.session_state.step = 1
        st.rerun()

# STEP 3: Plan Display
elif st.session_state.step == 3:
    st.header(T[lang]["step3_title"])
    st.write(T[lang]["step3_desc"])
    
    # Handle Translation Options
    if lang == "Arabic":
        col_lang, _ = st.columns([1, 4])
        with col_lang:
            default_idx = 1 if st.session_state.view_language == "Arabic" else 0
            st.selectbox("View Language / عرض اللغة", options=[T[lang]["view_original"], T[lang]["view_translated"]], index=default_idx, key="view_language_sel_3")
            st.session_state.view_language = "Arabic" if st.session_state.view_language_sel_3 == T[lang]["view_translated"] else "English"
        
        if st.session_state.view_language == "Arabic" and st.session_state.plan_result_ar is None:
            with st.spinner(T[lang]["translate_btn"]):
                try:
                    st.session_state.plan_result_ar = translate_data_to_arabic(st.session_state.plan_result, "Workshop plan outline")
                except Exception as e:
                    st.error(f"{T[lang]['translate_failed']}: {e}")
                    st.session_state.view_language = "English"
    else:
        st.session_state.view_language = "English"
        
    active_plan = st.session_state.plan_result_ar if (st.session_state.view_language == "Arabic" and st.session_state.plan_result_ar) else st.session_state.plan_result
    
    if active_plan:
        # Chosen title
        st.subheader(f"Title: {st.session_state.chosen_title}")
        
        # Objectives Card
        st.markdown(f"### {T[lang]['objectives']}")
        with st.container():
            st.markdown("<div class='title-card'>", unsafe_allow_html=True)
            for obj in active_plan.get("learning_objectives", []):
                st.markdown(f"- **{obj}**")
            st.markdown("</div>", unsafe_allow_html=True)
            
        # Timeline Card (outline)
        st.markdown(f"### {T[lang]['outline']}")
        
        # Icon guessing logic
        role_icons = {
            "opening": "🏁",
            "explain": "📖",
            "lab": "🧪",
            "break": "☕",
            "qna": "💬",
            "competition": "🏆",
            "closing": "🏁"
        }
        
        for idx, item in enumerate(active_plan.get("outline", [])):
            role = item.get("role", "explain")
            icon = role_icons.get(role.lower(), "📖")
            
            with st.container():
                col_time, col_details = st.columns([1, 5])
                with col_time:
                    st.markdown(f"<div style='text-align: center; background-color: #161b22; padding: 12px 6px; border-radius: 8px; border: 1px solid #30363d; color: #58a6ff; font-weight: bold; font-size: 15px;'>⏱ {item.get('duration_minutes')} min</div>", unsafe_allow_html=True)
                with col_details:
                    st.markdown(f"<div class='timeline-card' style='margin-bottom: 12px; background-color: #161b22; border: 1px solid #30363d; border-radius: 8px; padding: 16px;'><strong style='font-size: 16px; color: #f0f6fc;'>{icon} {item.get('section')}</strong><p style='margin: 8px 0 0 0; font-size: 14px; color: #c9d1d9;'>{item.get('description')}</p></div>", unsafe_allow_html=True)
        
        st.markdown("---")
        
        # Refinement Loop
        st.subheader(T[lang]["refine_plan"])
        feedback = st.text_area(T[lang]["refine_plan"], placeholder=T[lang]["refine_plan_placeholder"], label_visibility="collapsed")
        
        if st.button(T[lang]["btn_refine"]):
            if feedback.strip():
                with st.spinner("Revising plan outline..." if lang == "English" else "جاري مراجعة المخطط..."):
                    try:
                        context = f"Plan structure for title='{st.session_state.chosen_title}'"
                        revised = revise_output(st.session_state.plan_result, feedback, context)
                        st.session_state.plan_result = revised
                        st.session_state.plan_result_ar = None # invalidate cache
                        st.success("Plan revised successfully!" if lang == "English" else "تم مراجعة المخطط بنجاح!")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Failed to revise plan: {e}")
                        
        st.markdown("---")
        
        # Generate Slide Content Option Form
        st.markdown("---")
        st.subheader(T[lang]["slide_options"])
        st.radio(
            T[lang]["gen_mode_label"],
            options=["ai", "mimic"],
            format_func=lambda x: T[lang]["mode_ai_choice"] if x == "ai" else T[lang]["mode_mimic"],
            key="slide_gen_mode_choice"
        )
        
        mimic_file = None
        if st.session_state.slide_gen_mode_choice == "mimic":
            mimic_file = st.file_uploader(
                T[lang]["upload_label"],
                type=["json", "txt", "pptx"],
                key="slide_ref_file"
            )
            
        st.markdown("<div style='margin-bottom: 20px;'></div>", unsafe_allow_html=True)
        
        # Generate Slide Content
        col_nav1, col_nav2 = st.columns(2)
        with col_nav1:
            if st.button(T[lang]["btn_back"], use_container_width=True):
                st.session_state.step = 2
                st.rerun()
        with col_nav2:
            if st.button(T[lang]["btn_generate_content"], type="primary", use_container_width=True):
                # Verify mimic input
                mimic_example_text = None
                if st.session_state.slide_gen_mode_choice == "mimic":
                    if mimic_file is None:
                        st.error("❌ Please upload a reference slide file (.json, .txt, or .pptx) first." if lang == "English" else "❌ يرجى رفع ملف مرجع الشرائح (.json أو .txt أو .pptx) أولاً.")
                        st.stop()
                    else:
                        try:
                            if mimic_file.name.endswith(".pptx"):
                                # Extract slide content text from PPTX
                                mimic_example_text = extract_pptx_slides_text(mimic_file)
                            elif mimic_file.name.endswith(".json"):
                                file_bytes = mimic_file.read()
                                try:
                                    parsed_json = json.loads(file_bytes.decode("utf-8"))
                                    mimic_example_text = json.dumps(parsed_json, indent=2, ensure_ascii=False)
                                except Exception:
                                    mimic_example_text = file_bytes.decode("utf-8")
                            else:
                                file_bytes = mimic_file.read()
                                mimic_example_text = file_bytes.decode("utf-8")
                        except Exception as e:
                            st.error(f"Failed to read file: {e}")
                            st.stop()
                            
                # Trigger content generation with progress updates section-by-section
                progress_bar = st.progress(0.0)
                status_text = st.empty()
                
                def update_progress(index, total, section_name):
                    progress_val = float(index) / float(total)
                    progress_bar.progress(progress_val)
                    status_text.markdown(
                        f"⚙️ **Generating Slide Content ({index+1}/{total})**: {section_name}..."
                        if lang == "English" else
                        f"⚙️ **جاري توليد محتوى الشرائح ({index+1}/{total})**: {section_name}..."
                    )
                
                try:
                    section_styles = {}
                    default_style = st.session_state.get("slide_style", "Clean & Minimal")
                    if st.session_state.get("style_application") in ["Choose per section", "اختر لكل قسم"]:
                        for idx_sec, sec in enumerate(st.session_state.plan_result["outline"]):
                            style_val = st.session_state.get(f"style_section_{idx_sec}", default_style)
                            section_styles[sec["section"]] = style_val
                            
                    content = generate_content_with_sources(
                        title=st.session_state.chosen_title,
                        learning_objectives=st.session_state.plan_result["learning_objectives"],
                        outline=st.session_state.plan_result["outline"],
                        mimic_example=mimic_example_text,
                        default_style=default_style,
                        section_styles=section_styles,
                        progress_callback=update_progress
                    )
                    
                    st.session_state.content_result = content
                    st.session_state.step = 4
                    clear_block_keys() # Reset editor state keys
                    
                    progress_bar.empty()
                    status_text.empty()
                    st.rerun()
                except Exception as e:
                    progress_bar.empty()
                    status_text.empty()
                    st.error(f"Error during slide content generation: {e}")

# STEP 4: Content with Manual Editing
elif st.session_state.step == 4:
    st.header(T[lang]["step4_title"])
    st.write(T[lang]["step4_desc"])
    
    # Handle Translation Options
    if lang == "Arabic":
        col_lang, _ = st.columns([1, 4])
        with col_lang:
            default_idx = 1 if st.session_state.view_language == "Arabic" else 0
            st.selectbox("View Language / عرض اللغة", options=[T[lang]["view_original"], T[lang]["view_translated"]], index=default_idx, key="view_language_sel_4")
            st.session_state.view_language = "Arabic" if st.session_state.view_language_sel_4 == T[lang]["view_translated"] else "English"
        
        if st.session_state.view_language == "Arabic" and st.session_state.content_result_ar is None:
            with st.spinner(T[lang]["translate_btn"]):
                try:
                    st.session_state.content_result_ar = translate_data_to_arabic(st.session_state.content_result, "Workshop slides content")
                except Exception as e:
                    st.error(f"{T[lang]['translate_failed']}: {e}")
                    st.session_state.view_language = "English"
    else:
        st.session_state.view_language = "English"
        
    active_content = st.session_state.content_result_ar if (st.session_state.view_language == "Arabic" and st.session_state.content_result_ar) else st.session_state.content_result
    
    if active_content and "slides" in active_content:
        slides_list = active_content.get("slides", [])

        # ─────────────────────────────────────────────
        # GLOBAL AI EDIT — edits the entire slide deck (can add, remove, rewrite)
        # ─────────────────────────────────────────────
        with st.expander("✨ " + ("Edit Whole Presentation with AI" if lang == "English" else "تعديل كل العرض بالذكاء الاصطناعي"), expanded=False):
            st.caption("Instruct the AI to rewrite, add, remove, or reorganize slides." if lang == "English" else "اطلب من الذكاء الاصطناعي إعادة كتابة الشرائح، إضافتها، إزالتها، أو إعادة تنظيمها.")
            global_edit_instruction = st.text_area(
                "Instruction for the whole presentation:" if lang == "English" else "التعليمة لكامل العرض:",
                placeholder="e.g. Add 3 more slides about advanced techniques. / Make all slides shorter. / Change the tone to be more academic." if lang == "English" else "مثال: أضف 3 شرائح إضافية حول التقنيات المتقدمة. / اجعل جميع الشرائح أقصر.",
                key="global_ai_edit_instruction",
                height=90
            )
            if st.button("🚀 " + ("Apply to Presentation" if lang == "English" else "طبّق على العرض"), key="global_ai_edit_btn", type="primary", use_container_width=True):
                instruction = st.session_state.get("global_ai_edit_instruction", "").strip()
                if not instruction:
                    st.warning("Please enter an instruction first." if lang == "English" else "الرجاء إدخال تعليمة أولاً.")
                else:
                    from llm_client import ask_llm_for_json
                    with st.spinner("Analyzing slide deck and applying changes..." if lang == "English" else "جاري تحليل العرض وتطبيق التغييرات..."):
                        try:
                            edit_prompt = f"""You are an expert presentation editor. The user wants to modify their presentation slides.

INSTRUCTION:
"{instruction}"

CURRENT PRESENTATION SLIDES (JSON):
{json.dumps(active_content, ensure_ascii=False, indent=2)}

You can rewrite slide contents, add new slides, delete slides, or change their order.
Return the complete updated presentation as JSON with the exact same format:
{{
  "slides": [
    {{
      "slide_number": 1,
      "slide_title": "...",
      "content_type": "content_slide" or "columns_3_slide" or "timeline_slide",
      "section": "...",
      "blocks": [
         // list of blocks: heading, paragraph, bullet_list, activity, image_placeholder
      ],
      "sources": [
         // optional list of citation source dicts: author, year, title, url, exact_quote
      ],
      "slide_style": "Clean & Minimal" or "Bold & Impactful" or "Visual & Diagram-heavy" or "Data & Research" or "Interactive & Workshop"
    }},
    ...
  ]
}}

Ensure every slide you output has valid block structures following the schema:
- content_slide / roadmap_slide: block types: heading, paragraph, bullet_list, activity, image_placeholder
- columns_3_slide: block types: paragraph, columns_3 (containing a "columns" list with "heading", "text", "image_placeholder")
- timeline_slide: block types: timeline (containing "events" list with "title", "date", "text"), image_placeholder

Do NOT include any extra keys or text other than the JSON."""
                            edited_deck = ask_llm_for_json(edit_prompt)
                            if not edited_deck or "slides" not in edited_deck or not isinstance(edited_deck["slides"], list):
                                st.error("Failed to apply edits. AI returned an invalid schema.")
                            else:
                                # Re-sequence slide numbers
                                for g_idx, g_slide in enumerate(edited_deck["slides"], 1):
                                    g_slide["slide_number"] = g_idx
                                    if "slide_style" not in g_slide:
                                        g_slide["slide_style"] = active_content["slides"][0].get("slide_style", "Clean & Minimal") if active_content["slides"] else "Clean & Minimal"
                                    if "sources" not in g_slide:
                                        g_slide["sources"] = []
                                
                                if st.session_state.view_language == "Arabic" and st.session_state.content_result_ar:
                                    st.session_state.content_result_ar = edited_deck
                                else:
                                    st.session_state.content_result = edited_deck
                                st.success("✅ " + ("Presentation updated successfully!" if lang == "English" else "تم تحديث العرض بنجاح!"))
                                st.rerun()
                        except Exception as e_g:
                            st.error(f"Error applying AI edits: {e_g}")

        st.markdown("---")


        for slide_idx, slide in enumerate(slides_list):
            with st.container():
                # Slide card with border
                st.markdown(f"""
                <div style="
                    border: 2px solid #667eea;
                    border-radius: 12px;
                    padding: 20px;
                    margin-bottom: 30px;
                    background: white;
                    box-shadow: 0 4px 6px rgba(0,0,0,0.1);
                    color: #333;
                ">
                    <div style="display: flex; justify-content: space-between; align-items: center; border-bottom: 2px solid #667eea; padding-bottom: 10px; margin-bottom: 15px;">
                        <h3 style="color: #667eea; margin: 0; font-family: sans-serif;">📄 Slide {slide['slide_number']}: {slide['slide_title']}</h3>
                        <span style="background: #667eea; color: white; padding: 4px 12px; border-radius: 20px; font-size: 12px; font-weight: bold; font-family: sans-serif;">
                            {slide.get('content_type', 'content_slide')}
                        </span>
                    </div>
                """, unsafe_allow_html=True)
                
                # Render each block with spacing
                for b_idx, block in enumerate(slide.get("blocks", [])):
                    b_type = block.get("type")
                    if b_type == "heading":
                        st.markdown(f"<h4 style='color: #333; margin: 10px 0; font-family: sans-serif;'>{block.get('text', '')}</h4>", unsafe_allow_html=True)
                        
                    elif b_type == "paragraph":
                        # Editable paragraph
                        new_text = st.text_area(
                            f"",
                            value=block.get("text", ""),
                            key=f"edit_para_{slide_idx}_{b_idx}",
                            height=75,
                            label_visibility="collapsed",
                            placeholder="One sentence (10-15 words) max..."
                        )
                        if new_text != block.get("text"):
                            block["text"] = new_text
                            if st.session_state.view_language == "Arabic":
                                st.session_state.content_result_ar["slides"][slide_idx]["blocks"][b_idx]["text"] = new_text
                            else:
                                st.session_state.content_result["slides"][slide_idx]["blocks"][b_idx]["text"] = new_text
                            
                    elif b_type == "bullet_list":
                        # Editable bullet list
                        items_list = block.get("items", [])
                        items_text = "\n".join(items_list)
                        new_items = st.text_area(
                            f"",
                            value=items_text,
                            key=f"edit_bullets_{slide_idx}_{b_idx}",
                            height=125,
                            label_visibility="collapsed",
                            placeholder="One bullet per line (3-5 words each)"
                        )
                        if new_items != items_text:
                            parsed_items = [item.strip() for item in new_items.split("\n") if item.strip()]
                            block["items"] = parsed_items
                            if st.session_state.view_language == "Arabic":
                                st.session_state.content_result_ar["slides"][slide_idx]["blocks"][b_idx]["items"] = parsed_items
                            else:
                                st.session_state.content_result["slides"][slide_idx]["blocks"][b_idx]["items"] = parsed_items
                            
                    elif b_type == "image_placeholder":
                        st.markdown(f"""
                        <div style="margin: 10px 0; padding: 20px; background: #f0f4ff; border: 2px dashed #667eea; border-radius: 8px; text-align: center; color: #333;">
                            <p style="margin: 0; color: #666; font-family: sans-serif;">🖼️ {block.get('text', 'Image placeholder')}</p>
                            <p style="margin: 5px 0 0 0; font-size: 12px; color: #999; font-family: sans-serif;">(Add your own image)</p>
                        </div>
                        """, unsafe_allow_html=True)
                        
                    elif b_type == "roadmap":
                        # Editable roadmap items
                        st.markdown("**Roadmap Modules:**")
                        items_list = block.get("items", [])
                        for i_idx, item in enumerate(items_list):
                            c_col1, c_col2 = st.columns(2)
                            with c_col1:
                                new_title = st.text_input(
                                    f"Block Title {i_idx+1}",
                                    value=item.get("title", ""),
                                    key=f"roadmap_title_{slide_idx}_{b_idx}_{i_idx}"
                                )
                                if new_title != item.get("title"):
                                    item["title"] = new_title
                            with c_col2:
                                new_detail = st.text_input(
                                    f"Detail {i_idx+1}",
                                    value=item.get("detail", ""),
                                    key=f"roadmap_detail_{slide_idx}_{b_idx}_{i_idx}"
                                )
                                if new_detail != item.get("detail"):
                                    item["detail"] = new_detail
                                    
                    elif b_type == "columns_3":
                        # Editable 3 columns
                        st.markdown("**Columns Block (Max 3):**")
                        cols_list = block.get("columns", [])
                        for c_idx, col in enumerate(cols_list):
                            with st.expander(f"Column {c_idx+1}: {col.get('heading', 'Untitled')}", expanded=True):
                                new_head = st.text_input(
                                    "Column Heading",
                                    value=col.get("heading", ""),
                                    key=f"col_head_{slide_idx}_{b_idx}_{c_idx}"
                                )
                                if new_head != col.get("heading"):
                                    col["heading"] = new_head
                                new_text = st.text_area(
                                    "Column Text",
                                    value=col.get("text", ""),
                                    key=f"col_text_{slide_idx}_{b_idx}_{c_idx}",
                                    height=60
                                )
                                if new_text != col.get("text"):
                                    col["text"] = new_text
                                new_img = st.text_input(
                                    "Column Image Concept",
                                    value=col.get("image_placeholder", ""),
                                    key=f"col_img_{slide_idx}_{b_idx}_{c_idx}"
                                )
                                if new_img != col.get("image_placeholder"):
                                    col["image_placeholder"] = new_img
                                    
                    elif b_type == "timeline":
                        # Editable vertical timeline
                        st.markdown("**Timeline Eras/Events:**")
                        events_list = block.get("events", [])
                        for e_idx, event in enumerate(events_list):
                            with st.expander(f"Event {e_idx+1}: {event.get('title', 'Untitled')}", expanded=True):
                                col_t1, col_t2 = st.columns(2)
                                with col_t1:
                                    new_evt_title = st.text_input(
                                        "Event Title",
                                        value=event.get("title", ""),
                                        key=f"evt_title_{slide_idx}_{b_idx}_{e_idx}"
                                    )
                                    if new_evt_title != event.get("title"):
                                        event["title"] = new_evt_title
                                with col_t2:
                                    new_evt_date = st.text_input(
                                        "Event Date/Era",
                                        value=event.get("date", ""),
                                        key=f"evt_date_{slide_idx}_{b_idx}_{e_idx}"
                                    )
                                    if new_evt_date != event.get("date"):
                                        event["date"] = new_evt_date
                                new_evt_text = st.text_area(
                                    "Event Details",
                                    value=event.get("text", ""),
                                    key=f"evt_text_{slide_idx}_{b_idx}_{e_idx}",
                                    height=60
                                )
                                if new_evt_text != event.get("text"):
                                    event["text"] = new_evt_text
                                    
                    elif b_type == "activity":
                        # Editable activity block
                        st.markdown("**⚡ Activity / Callout:**")
                        new_act = st.text_area(
                            "",
                            value=block.get("text", ""),
                            key=f"edit_activity_{slide_idx}_{b_idx}",
                            height=75,
                            label_visibility="collapsed",
                            placeholder="Activity description..."
                        )
                        if new_act != block.get("text"):
                            block["text"] = new_act
                            if st.session_state.view_language == "Arabic":
                                st.session_state.content_result_ar["slides"][slide_idx]["blocks"][b_idx]["text"] = new_act
                            else:
                                st.session_state.content_result["slides"][slide_idx]["blocks"][b_idx]["text"] = new_act
                
                # Sources section
                if slide.get("sources"):
                    with st.expander("📚 Sources", expanded=True):
                        for source in slide["sources"]:
                            # Guard: LLM occasionally returns a plain string instead of a dict
                            if not isinstance(source, dict):
                                st.markdown(f"<div style='padding:8px;background:#f5f5f5;border-radius:4px;margin-bottom:8px;border-left:3px solid #4CAF50;color:#333;font-size:13px;font-family:sans-serif;'>{source}</div>", unsafe_allow_html=True)
                                continue
                            st.markdown(f"""
                            <div style="padding: 8px; background: #f5f5f5; border-radius: 4px; margin-bottom: 8px; border-left: 3px solid #4CAF50; color: #333;">
                                <p style="margin: 0; font-size: 13px; font-family: sans-serif;">
                                    <strong>{source.get('author', 'Unknown')} ({source.get('year', 'n.d.')})</strong> - {source.get('title', 'Untitled')}
                                </p>
                                <p style="margin: 2px 0 0 0; font-size: 12px; color: #667eea; font-family: sans-serif;">
                                    🔗 <a href="{source.get('url', '#')}" target="_blank" style="color: #667eea; text-decoration: underline;">{source.get('url', 'N/A')}</a>
                                </p>
                                <p style="margin: 4px 0 0 0; font-size: 12px; color: #666; font-style: italic; font-family: sans-serif;">
                                    "{source.get('exact_quote', 'No quote available')}"
                                </p>
                            </div>
                            """, unsafe_allow_html=True)
                
                st.markdown("</div>", unsafe_allow_html=True)
                
                # ✨ Edit with AI per slide
                with st.expander(f"✨ {'Edit Slide with AI' if lang == 'English' else 'تعديل الشريحة بالذكاء الاصطناعي'}", expanded=False):
                    ai_instruction = st.text_input(
                        "Tell AI how to edit this slide:" if lang == "English" else "أخبر الذكاء الاصطناعي كيف يعدّل هذه الشريحة:",
                        placeholder="E.g., Make it simpler, add more data, make it more engaging..." if lang == "English" else "مثال: اجعلها أبسط، أضف بيانات أكثر، اجعلها أكثر تفاعلية...",
                        key=f"ai_edit_instruction_{slide_idx}"
                    )
                    if st.button(f"✨ {'Apply AI Edit' if lang == 'English' else 'تطبيق التعديل'}", key=f"ai_edit_btn_{slide_idx}", use_container_width=True):
                        if ai_instruction and ai_instruction.strip():
                            with st.spinner("AI is rewriting this slide..." if lang == "English" else "الذكاء الاصطناعي يعيد كتابة الشريحة..."):
                                try:
                                    from llm_client import ask_llm_for_json
                                    import json as _json
                                    current_slide_json = _json.dumps(slide, ensure_ascii=False, indent=2)
                                    edit_prompt = f"""You are editing a SINGLE presentation slide. The user wants you to modify it.

USER INSTRUCTION: {ai_instruction}

CURRENT SLIDE JSON:
{current_slide_json}

RULES:
1. Keep the SAME JSON structure (same keys: slide_number, section, slide_title, content_type, blocks, sources, speaker_notes)
2. Keep slide_number, section, and content_type UNCHANGED
3. Apply the user's instruction to the text content (slide_title, blocks text, bullet items, speaker_notes)
4. Each paragraph = 1 sentence (10-15 words max)
5. Each bullet item = 3-5 words
6. Keep the same number and types of blocks unless the instruction specifically asks to add/remove
7. Keep all sources unchanged unless the instruction specifically mentions sources

Return ONLY the modified slide as a JSON object (NOT wrapped in an array).
"""
                                    edited_slide = ask_llm_for_json(edit_prompt)
                                    
                                    # Preserve immutable fields
                                    edited_slide["slide_number"] = slide["slide_number"]
                                    edited_slide["section"] = slide["section"]
                                    edited_slide["content_type"] = slide.get("content_type", "content_slide")
                                    edited_slide["slide_style"] = slide.get("slide_style", "Clean & Minimal")
                                    if "sources" not in edited_slide:
                                        edited_slide["sources"] = slide.get("sources", [])
                                    
                                    # Update the slide in the active content
                                    if st.session_state.view_language == "Arabic" and st.session_state.content_result_ar:
                                        st.session_state.content_result_ar["slides"][slide_idx] = edited_slide
                                    else:
                                        st.session_state.content_result["slides"][slide_idx] = edited_slide
                                    
                                    clear_block_keys()
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"{'AI edit failed' if lang == 'English' else 'فشل التعديل بالذكاء الاصطناعي'}: {e}")
                        else:
                            st.warning("Please type an instruction first." if lang == "English" else "يرجى كتابة تعليمات أولاً.")
                
        # PowerPoint Export Settings
        st.markdown("---")
        st.subheader("PowerPoint Export Options 📊" if lang == "English" else "خيارات تصدير البوربوينت 📊")
        
        pptx_options = [
            "Generate default PowerPoint (Clean, Bold, etc.)" if lang == "English" else "توليد ملف بوربوينت افتراضي",
            "Upload a custom template to mimic its style and colors" if lang == "English" else "رفع قالب بوربوينت مخصص لمحاكاة النمط والألوان"
        ]
        
        st.radio(
            "Select PowerPoint Generation Mode:" if lang == "English" else "اختر طريقة توليد البوربوينت:",
            options=pptx_options,
            key="pptx_generation_mode"
        )
        
        st.checkbox(
            "Include Stock Photos in Slides" if lang == "English" else "تضمين الصور في الشرائح",
            value=True,
            key="pptx_include_images"
        )
        
        if st.session_state.pptx_generation_mode == pptx_options[0]:
            st.selectbox(
                "Select default presentation style" if lang == "English" else "اختر نمط العرض التقديمي الافتراضي",
                options=[
                    "Clean & Minimal", 
                    "Bold & Impactful", 
                    "Visual & Diagram-heavy",
                    "Data & Research",
                    "Interactive & Workshop"
                ],
                key="slide_style",
            )
            
            st.radio(
                "Apply style to:" if lang == "English" else "تطبيق النمط على:",
                options=["All slides", "Choose per section"] if lang == "English" else ["جميع الشرائح", "اختر لكل قسم"],
                key="style_application",
                horizontal=True
            )
            
            if st.session_state.style_application in ["Choose per section", "اختر لكل قسم"] and "plan_result" in st.session_state and st.session_state.plan_result:
                outline = st.session_state.plan_result.get("outline", [])
                for i, section in enumerate(outline):
                    st.selectbox(
                        f"Style for {section['section']}" if lang == "English" else f"نمط قسم: {section['section']}",
                        options=["Clean & Minimal", "Bold & Impactful", "Visual & Diagram-heavy", "Data & Research", "Interactive & Workshop"],
                        key=f"style_section_{i}"
                    )
        
        elif st.session_state.pptx_generation_mode == pptx_options[1]:
            st.file_uploader(
                "Upload a PowerPoint presentation (.pptx) to mimic:" if lang == "English" else "ارفع ملف عرض تقديمي (.pptx) لمحاكاته:",
                type=["pptx"],
                key="pptx_custom_template_file"
            )

        st.markdown("---")
        
        # Navigation
        col1, col2, col3 = st.columns([1, 1, 1])
        with col1:
            if st.button("← Back to Plan" if lang == "English" else "← العودة إلى الخطة", use_container_width=True):
                st.session_state.step = 3
                st.rerun()
        with col2:
            if st.button("🧪 Generate Labs →" if lang == "English" else "🧪 توليد التدريبات العملي →", type="primary", use_container_width=True):
                # Pre-extract lab contexts for step 5
                try:
                    contexts = extract_lab_contexts(st.session_state.plan_result["outline"])
                    st.session_state.lab_contexts = contexts
                    st.session_state.step = 5
                    st.rerun()
                except Exception as e:
                    st.error(f"Failed to extract lab contexts: {e}")
        with col3:
            if st.button("📊 Generate PowerPoint" if lang == "English" else "📊 توليد ملف بوربوينت", use_container_width=True):
                with st.spinner("Creating PowerPoint..."):
                    try:
                        from pptx_export import export_to_pptx
                        import os
                        
                        pptx_opts = [
                            "Generate default PowerPoint (Clean, Bold, etc.)" if lang == "English" else "توليد ملف بوربوينت افتراضي",
                            "Upload a custom template to mimic its style and colors" if lang == "English" else "رفع قالب بوربوينت مخصص لمحاكاة النمط والألوان"
                        ]
                        
                        template_file_path = None
                        if st.session_state.get("pptx_generation_mode") == pptx_opts[1] and st.session_state.get("pptx_custom_template_file") is not None:
                            uploaded = st.session_state.pptx_custom_template_file
                            temp_path = "temp_uploaded_template.pptx"
                            with open(temp_path, "wb") as tf_file:
                                tf_file.write(uploaded.getbuffer())
                            template_file_path = temp_path
                            
                        # Build a copy of slides with chosen styles
                        import copy
                        slides_to_export = copy.deepcopy(active_content["slides"])
                        if st.session_state.get("pptx_generation_mode") == pptx_opts[0]:
                            chosen_style = st.session_state.get("slide_style", "Clean & Minimal")
                            style_app = st.session_state.get("style_application", "All slides")
                            
                            for idx, slide_data in enumerate(slides_to_export):
                                if style_app in ["All slides", "جميع الشرائح"]:
                                    slide_data["slide_style"] = chosen_style
                                else:
                                    section_name = slide_data.get("section", "")
                                    outline = st.session_state.plan_result.get("outline", [])
                                    section_idx = 0
                                    for o_idx, o_sec in enumerate(outline):
                                        if o_sec.get("section", "") == section_name:
                                            section_idx = o_idx
                                            break
                                    sec_style = st.session_state.get(f"style_section_{section_idx}", chosen_style)
                                    slide_data["slide_style"] = sec_style

                        pptx_path = export_to_pptx(
                            title=st.session_state.chosen_title,
                            slides=slides_to_export,
                            template_path=template_file_path,
                            include_images=st.session_state.get("pptx_include_images", True)
                        )
                        
                        if template_file_path and os.path.exists(template_file_path):
                            try:
                                os.remove(template_file_path)
                            except Exception:
                                pass
                        with open(pptx_path, "rb") as f:
                            st.download_button(
                                label="⬇️ Download PowerPoint" if lang == "English" else "⬇️ تحميل ملف البوربوينت",
                                data=f,
                                file_name=f"{st.session_state.chosen_title.replace(' ', '_')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True
                            )
                    except Exception as e:
                        st.error(f"❌ Error: {str(e)}")

# STEP 5: Lab Generator
elif st.session_state.step == 5:
    st.header(T[lang]["step5_title"])
    st.write(T[lang]["step5_desc"])
    
    # Handle Translation Options
    if lang == "Arabic":
        col_lang, _ = st.columns([1, 4])
        with col_lang:
            default_idx = 1 if st.session_state.view_language == "Arabic" else 0
            st.selectbox("View Language / عرض اللغة", options=[T[lang]["view_original"], T[lang]["view_translated"]], index=default_idx, key="view_language_sel_5")
            st.session_state.view_language = "Arabic" if st.session_state.view_language_sel_5 == T[lang]["view_translated"] else "English"
        
        if st.session_state.view_language == "Arabic" and st.session_state.labs_result is not None and st.session_state.labs_result_ar is None:
            with st.spinner(T[lang]["translate_btn"]):
                try:
                    st.session_state.labs_result_ar = translate_data_to_arabic(st.session_state.labs_result, "Workshop hands-on labs and keys")
                except Exception as e:
                    st.error(f"{T[lang]['translate_failed']}: {e}")
                    st.session_state.view_language = "English"
    else:
        st.session_state.view_language = "English"

    lab_contexts = st.session_state.lab_contexts
    
    if not lab_contexts:
        st.info("No hands-on labs required for this workshop outline." if lang == "English" else "لا توجد مختبرات تطبيقية مطلوبة في هذا الهيكل لورشة العمل.")
        col_nav1, col_nav2 = st.columns(2)
        with col_nav1:
            if st.button(T[lang]["btn_back"], use_container_width=True):
                st.session_state.step = 4
                st.rerun()
        with col_nav2:
            if st.button(T[lang]["btn_next"], type="primary", use_container_width=True):
                st.session_state.step = 6
                st.rerun()
                
    else:
        # Load suggestions if empty
        if st.session_state.lab_suggestions is None:
            with st.spinner("AI is analyzing slide content to suggest optimal lab designs..." if lang == "English" else "جاري تحليل المحتوى واقتراح الأفكار للمختبرات..."):
                try:
                    sugg = suggest_lab_types(st.session_state.chosen_title, lab_contexts)
                    st.session_state.lab_suggestions = sugg
                    
                    # Pre-fill confirmed types
                    for s in sugg.get("suggestions", []):
                        idx = s["outline_index"]
                        if idx not in st.session_state.confirmed_lab_types:
                            st.session_state.confirmed_lab_types[idx] = s["lab_type"]
                except Exception as e:
                    st.error(f"Failed to fetch lab type suggestions: {e}")
                    
        # Configuration list
        st.subheader("Lab Settings & Confirmation" if lang == "English" else "إعدادات وتأكيد المختبرات")
        
        # List configs for each lab slot
        for idx, ctx in enumerate(lab_contexts):
            out_idx = ctx["outline_index"]
            covers_text = ", ".join(ctx["covers_sections"])
            
            # Find matching suggestion if exists
            sugg_item = None
            if st.session_state.lab_suggestions:
                for s in st.session_state.lab_suggestions.get("suggestions", []):
                    if s["outline_index"] == out_idx:
                        sugg_item = s
                        break
                        
            st.markdown(f"**🔬 {T[lang]['lab_slot']} #{idx+1} ({ctx['duration_minutes']} min)**")
            st.write(f"*{T[lang]['covers']}: {covers_text}*")
            
            if sugg_item:
                st.info(f"💡 {T[lang]['suggested_type']}: **{sugg_item['lab_type']}** — *{sugg_item['reason']}*")
                
            col_type, col_notes = st.columns([1, 2])
            with col_type:
                # Pre-fill with suggestion
                default_type = st.session_state.confirmed_lab_types.get(out_idx, "coding")
                st.selectbox(
                    T[lang]["confirm_type"],
                    options=["coding", "conceptual"],
                    index=0 if default_type == "coding" else 1,
                    key=f"lab_type_sel_{out_idx}"
                )
                st.session_state.confirmed_lab_types[out_idx] = st.session_state[f"lab_type_sel_{out_idx}"]
                
            with col_notes:
                st.text_input(
                    T[lang]["custom_notes"],
                    placeholder="E.g., Keep it basic, explain loops...",
                    key=f"lab_notes_{out_idx}"
                )
                st.session_state.lab_customizations[out_idx] = st.session_state[f"lab_notes_{out_idx}"]
                
            st.markdown("---")
            
        # Trigger labs generation
        if st.button(T[lang]["btn_generate_labs"], type="primary", use_container_width=True):
            labs = []
            progress_bar = st.progress(0.0)
            status_text = st.empty()
            
            error_occurred = False
            for idx, ctx in enumerate(lab_contexts):
                out_idx = ctx["outline_index"]
                l_type = st.session_state.confirmed_lab_types.get(out_idx, "coding")
                l_notes = st.session_state.lab_customizations.get(out_idx, "")
                
                covers_str = ", ".join(ctx["covers_sections"])
                status_text.markdown(f"⚙️ **Generating Lab {idx+1}/{len(lab_contexts)}**: {covers_str}..." if lang == "English" else f"⚙️ **جاري توليد المختبر {idx+1}/{len(lab_contexts)}**: {covers_str}...")
                
                try:
                    # Always generate in English first
                    lab = generate_one_lab(
                        title=st.session_state.chosen_title,
                        lab_context=ctx,
                        lab_type=l_type,
                        content=st.session_state.content_result,
                        customization_notes=l_notes
                    )
                    labs.append(lab)
                except Exception as e:
                    st.error(f"Failed to generate lab {idx+1}: {e}")
                    error_occurred = True
                    break
                    
                progress_bar.progress((idx + 1) / len(lab_contexts))
                
            status_text.empty()
            progress_bar.empty()
            
            if not error_occurred:
                st.session_state.labs_result = {"labs": labs}
                st.session_state.labs_result_ar = None # reset translations cache
                
                # Automatically save notebooks to disk
                try:
                    save_all_lab_notebooks(st.session_state.labs_result, output_dir="generated_labs")
                except Exception as e:
                    st.warning(f"Note: Could not build physical .ipynb files: {e}")
                    
                st.rerun()

        # Display generated labs if available
        active_labs = st.session_state.labs_result_ar if (st.session_state.view_language == "Arabic" and st.session_state.labs_result_ar) else st.session_state.labs_result
        
        # Helper function to load notebooks
        def get_notebook_download_data(lab: dict):
            from notebook_builder import sanitize_filename
            base_name = sanitize_filename(lab.get("title", "lab"))
            t_path = os.path.join("generated_labs", f"{base_name}_trainee.ipynb")
            s_path = os.path.join("generated_labs", f"{base_name}_solution.ipynb")
            
            t_data, s_data = None, None
            if os.path.exists(t_path):
                with open(t_path, "rb") as f:
                    t_data = f.read()
            if os.path.exists(s_path):
                with open(s_path, "rb") as f:
                    s_data = f.read()
            return t_data, s_data, f"{base_name}_trainee.ipynb", f"{base_name}_solution.ipynb"
            
        if active_labs and "labs" in active_labs:
            st.markdown(f"### 📋 {T[lang]['lab_materials']}")
            
            for idx, lab in enumerate(active_labs["labs"]):
                l_title = lab.get("title", "Hands-on Exercise")
                l_type = lab.get("lab_type", "coding")
                l_dur = lab.get("duration_minutes", 0)
                l_covers = ", ".join(lab.get("covers_sections", []))
                
                with st.expander(f"🔬 Lab {idx+1}: {l_title} ({l_type.upper()} | {l_dur} min)"):
                    st.write(f"*{T[lang]['covers']}: {l_covers}*")
                    
                    # Group into tabs to prevent chaos
                    tab_info, tab_trainee, tab_solution, tab_platforms = st.tabs([
                        "📖 Instructions & Notes" if lang == "English" else "📖 التعليمات والملاحظات",
                        "💻 Trainee Version" if lang == "English" else "💻 نسخة المتدرب",
                        "🔑 Solution Version" if lang == "English" else "🔑 نسخة الحل",
                        "🌐 Suggested Platforms" if lang == "English" else "🌐 المنصات المقترحة"
                    ])
                    
                    with tab_info:
                        st.markdown("#### Instructions")
                        st.markdown(lab.get("instructions", ""))
                        st.markdown(f"#### {T[lang]['instructor_notes']}")
                        st.info(lab.get("instructor_notes", ""))
                        
                    with tab_trainee:
                        t_data, s_data, t_name, s_name = get_notebook_download_data(lab)
                        if t_data and l_type == "coding":
                            st.download_button(
                                label="Download Trainee Notebook (.ipynb) 📥" if lang == "English" else "تحميل دفتر ملاحظات المتدرب (.ipynb) 📥",
                                data=t_data,
                                file_name=t_name,
                                mime="application/x-ipynb+json",
                                key=f"dl_trainee_{idx}"
                            )
                            st.markdown("---")
                            
                        if l_type == "coding":
                            for cell in lab.get("trainee_notebook_cells", []):
                                c_type = cell.get("cell_type", "markdown")
                                c_content = cell.get("content", "")
                                if c_type == "markdown":
                                    st.markdown(c_content)
                                else:
                                    st.code(c_content, language="python")
                        else:
                            for q_idx, q in enumerate(lab.get("questions", []), 1):
                                st.markdown(f"**Q{q_idx}: {q.get('question')}**")
                                
                    with tab_solution:
                        t_data, s_data, t_name, s_name = get_notebook_download_data(lab)
                        if s_data and l_type == "coding":
                            st.download_button(
                                label="Download Solution Notebook (.ipynb) 📥" if lang == "English" else "تحميل دفتر ملاحظات الحل (.ipynb) 📥",
                                data=s_data,
                                file_name=s_name,
                                mime="application/x-ipynb+json",
                                key=f"dl_solution_{idx}"
                            )
                            st.markdown("---")
                            
                        if l_type == "coding":
                            for cell in lab.get("solution_notebook_cells", []):
                                c_type = cell.get("cell_type", "markdown")
                                c_content = cell.get("content", "")
                                if c_type == "markdown":
                                    st.markdown(c_content)
                                else:
                                    st.code(c_content, language="python")
                        else:
                            for q_idx, q in enumerate(lab.get("questions", []), 1):
                                st.markdown(f"**Q{q_idx}: {q.get('question')}**")
                                st.success(f"Answer: {q.get('answer')}")
                                
                    with tab_platforms:
                        for plat in lab.get("suggested_platforms", []):
                            st.markdown(f"- **[{plat.get('name')}]({plat.get('url')})**: {plat.get('why_it_fits')}")
                        
            st.markdown("---")
            
            # AI Refinement loop for labs
            st.subheader(T[lang]["refine_labs"])
            feedback = st.text_area(T[lang]["refine_labs"], placeholder="E.g., Make the coding task a bit harder...", label_visibility="collapsed")
            
            if st.button(T[lang]["btn_refine"]):
                if feedback.strip():
                    with st.spinner("AI is revising the labs..." if lang == "English" else "جاري مراجعة المختبرات وتعديلها..."):
                        try:
                            context = f"Labs content generated for '{st.session_state.chosen_title}'"
                            # Always revise English source first
                            revised = revise_output(st.session_state.labs_result, feedback, context)
                            st.session_state.labs_result = revised
                            st.session_state.labs_result_ar = None # reset translations cache
                            st.success("Labs revised successfully!" if lang == "English" else "تم مراجعة المختبرات بنجاح!")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Failed to revise labs: {e}")
                            
        # Navigation
        col_nav1, col_nav2 = st.columns(2)
        with col_nav1:
            if st.button(T[lang]["btn_back"], use_container_width=True):
                st.session_state.step = 4
                st.rerun()
        with col_nav2:
            if st.button(T[lang]["btn_next"], type="primary", use_container_width=True):
                # Generate Quiz on continue
                with st.spinner("Generating end-of-workshop comprehensive quiz..." if lang == "English" else "جاري توليد اختبار نهاية ورشة العمل التقييمي..."):
                    try:
                        # Set default count to 6 if not already set
                        if "num_quiz_questions" not in st.session_state:
                            st.session_state.num_quiz_questions = 6
                        quiz = generate_quiz(
                            title=st.session_state.chosen_title,
                            outline=st.session_state.plan_result["outline"],
                            content=st.session_state.content_result,
                            min_questions=st.session_state.num_quiz_questions,
                            max_questions=st.session_state.num_quiz_questions
                        )
                        st.session_state.quiz_result = quiz
                        st.session_state.quiz_result_ar = None # Reset translation cache
                        st.session_state.quiz_approved = False
                        st.session_state.quiz_approved_snapshot = None
                        st.session_state.step = 6
                        st.rerun()
                    except Exception as e:
                        st.error(f"Error during quiz generation: {e}")

# STEP 6: Quiz Display & Export
elif st.session_state.step == 6:
    st.header(T[lang]["step6_title"])
    st.write(T[lang]["step6_desc"])
    
    # Handle Translation Options
    if lang == "Arabic":
        col_lang, _ = st.columns([1, 4])
        with col_lang:
            default_idx = 1 if st.session_state.view_language == "Arabic" else 0
            st.selectbox("View Language / عرض اللغة", options=[T[lang]["view_original"], T[lang]["view_translated"]], index=default_idx, key="view_language_sel_6")
            st.session_state.view_language = "Arabic" if st.session_state.view_language_sel_6 == T[lang]["view_translated"] else "English"
        
        if st.session_state.view_language == "Arabic" and st.session_state.quiz_result_ar is None:
            with st.spinner(T[lang]["translate_btn"]):
                try:
                    st.session_state.quiz_result_ar = translate_data_to_arabic(st.session_state.quiz_result, "Workshop final quiz")
                except Exception as e:
                    st.error(f"{T[lang]['translate_failed']}: {e}")
                    st.session_state.view_language = "English"
    else:
        st.session_state.view_language = "English"

    active_quiz_dict = st.session_state.quiz_result_ar if (st.session_state.view_language == "Arabic" and st.session_state.quiz_result_ar) else st.session_state.quiz_result
    
    if active_quiz_dict and "quiz" in active_quiz_dict:
        quiz_data = active_quiz_dict["quiz"]
        questions = quiz_data.get("questions", [])
        
        st.subheader(f"📝 {quiz_data.get('title', T[lang]['quiz_title'])} ({len(questions)} Questions)")
        
        # Quiz Settings & Regeneration
        st.markdown("### Quiz Settings & Regeneration 🔄" if lang == "English" else "### إعدادات الاختبار وإعادة التوليد 🔄")
        col_reg1, col_reg2 = st.columns([2, 1])
        with col_reg1:
            if "num_quiz_questions" not in st.session_state:
                st.session_state.num_quiz_questions = len(questions)
                
            num_q = st.slider(
                "Target Number of Questions" if lang == "English" else "العدد المستهدف للأسئلة",
                min_value=3,
                max_value=30,
                value=st.session_state.num_quiz_questions,
                step=3,
                key="step6_quiz_count_slider"
            )
            st.session_state.num_quiz_questions = num_q
        with col_reg2:
            st.markdown("<div style='margin-top: 28px;'></div>", unsafe_allow_html=True)
            if st.button("Regenerate Quiz 🔄" if lang == "English" else "إعادة توليد الاختبار 🔄", use_container_width=True, key="regenerate_quiz_btn"):
                with st.spinner("Generating new quiz..." if lang == "English" else "جاري توليد اختبار جديد..."):
                    try:
                        quiz = generate_quiz(
                            title=st.session_state.chosen_title,
                            outline=st.session_state.plan_result["outline"],
                            content=st.session_state.content_result,
                            min_questions=st.session_state.num_quiz_questions,
                            max_questions=st.session_state.num_quiz_questions
                        )
                        st.session_state.quiz_result = quiz
                        st.session_state.quiz_result_ar = None # reset cache
                        st.session_state.quiz_approved = False
                        st.session_state.quiz_approved_snapshot = None
                        st.success("Quiz regenerated!" if lang == "English" else "تم إعادة توليد الاختبار!")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Failed to generate: {e}")
        
        # Group questions by difficulty
        easy_q = [q for q in questions if q.get("difficulty") == "easy"]
        medium_q = [q for q in questions if q.get("difficulty") == "medium"]
        hard_q = [q for q in questions if q.get("difficulty") == "hard"]
        
        # Tabs for grouping
        tab_easy, tab_medium, tab_hard = st.tabs(["Easy 🟢", "Medium 🟡", "Hard 🔴"])
        
        with tab_easy:
            for i, q in enumerate(easy_q, 1):
                with st.expander(f"Question {i}: {q.get('question')}"):
                    for opt in q.get("options", []):
                        if opt == q.get("correct_answer"):
                            st.success(f"✓ {opt} ({T[lang]['correct_ans']})")
                        else:
                            st.info(opt)
                            
        with tab_medium:
            for i, q in enumerate(medium_q, 1):
                with st.expander(f"Question {i}: {q.get('question')}"):
                    for opt in q.get("options", []):
                        if opt == q.get("correct_answer"):
                            st.success(f"✓ {opt} ({T[lang]['correct_ans']})")
                        else:
                            st.info(opt)
                            
        with tab_hard:
            for i, q in enumerate(hard_q, 1):
                with st.expander(f"Question {i}: {q.get('question')}"):
                    for opt in q.get("options", []):
                        if opt == q.get("correct_answer"):
                            st.success(f"✓ {opt} ({T[lang]['correct_ans']})")
                        else:
                            st.info(opt)
                            
        st.markdown("---")

        # --- Approval gate. Everything below this point is ALWAYS built
        # from a FROZEN snapshot taken at the exact moment of approval —
        # never from the live, still-editable quiz_result. Every
        # regenerate/refine action above resets quiz_approved to False, so
        # a stale snapshot can never sit next to a newer, unapproved draft.
        st.markdown("### ✅ Final Approval" if lang == "English" else "### ✅ الاعتماد النهائي")

        if st.session_state.quiz_approved and st.session_state.quiz_approved_snapshot:
            st.success(
                "This quiz is approved. Every export below reflects this exact approved version."
                if lang == "English" else
                "تم اعتماد هذا الاختبار. كل ملفات التصدير أدناه تطابق هذه النسخة المعتمدة بالضبط."
            )
            if st.button(
                "🔓 Unlock to edit again"
                if lang == "English" else
                "🔓 فتح للتعديل",
                key="quiz_unlock_btn"
            ):
                st.session_state.quiz_approved = False
                st.session_state.quiz_approved_snapshot = None
                st.rerun()
        else:
            st.info(
                "Review the questions above. Exports (Word, Kahoot, JSON) only become available "
                "after you approve — and will always match this exact version, never a version "
                "you edit afterward."
                if lang == "English" else
                "راجعي الأسئلة أعلاه. التصدير (Word، كاهوت، JSON) ما يصير متاح إلا بعد الاعتماد، "
                "وديماً بيطابق هذي النسخة بالضبط، مو أي نسخة تعدّلينها بعدها."
            )
            if st.button(
                "✅ Approve Final Quiz" if lang == "English" else "✅ اعتماد الاختبار النهائي",
                type="primary",
                key="quiz_approve_btn"
            ):
                st.session_state.quiz_approved = True
                st.session_state.quiz_approved_snapshot = copy.deepcopy(active_quiz_dict)
                st.rerun()

        st.markdown("---")

        # Once the quiz is approved, Step 7 (the final quality checklist)
        # and the export downloads live on their own page — this keeps
        # step 6 focused on reviewing/approving the quiz, and step 7
        # focused on the final cross-check before anything is exportable.
        if st.session_state.quiz_approved and st.session_state.quiz_approved_snapshot:
            if st.button(
                "🔍 Continue to Quality Checklist (Step 7) ➔"
                if lang == "English" else
                "🔍 التالي: قائمة الجودة النهائية (الخطوة 7) ➔",
                type="primary",
                use_container_width=True,
                key="goto_step7_btn"
            ):
                st.session_state.step = 7
                st.rerun()

        
        st.markdown("---")
        
        # AI Refinement loop for quiz
        st.subheader(T[lang]["refine_quiz"])
        feedback = st.text_area(
            T[lang]["refine_quiz"],
            placeholder="E.g., Replace question 3 with a question about RAG agents...",
            label_visibility="collapsed",
            key="quiz_refine_feedback_input_text"
        )
        
        if st.button(T[lang]["btn_refine"], key="quiz_refine_apply_btn"):
            if feedback.strip():
                with st.spinner("AI is revising the quiz..." if lang == "English" else "جاري مراجعة وتعديل الأسئلة..."):
                    try:
                        context = f"Quiz questions generated for '{st.session_state.chosen_title}'"
                        # Always revise English source first
                        revised = revise_output(st.session_state.quiz_result, feedback, context)
                        st.session_state.quiz_result = revised
                        st.session_state.quiz_result_ar = None # reset translations cache
                        st.session_state.quiz_approved = False
                        st.session_state.quiz_approved_snapshot = None
                        st.success("Quiz revised successfully!" if lang == "English" else "تم مراجعة الاختبار بنجاح!")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Failed to revise quiz: {e}")
                        
    else:
        st.info("No quiz has been generated for this workshop yet. Select the number of questions and generate it below:" if lang == "English" else "لم يتم توليد اختبار لهذه الورشة بعد. اختر عدد الأسئلة وقم بالتوليد أدناه:")
        
        if "num_quiz_questions" not in st.session_state:
            st.session_state.num_quiz_questions = 6
            
        num_q = st.slider(
            "Target Number of Questions" if lang == "English" else "العدد المستهدف للأسئلة",
            min_value=3,
            max_value=30,
            value=st.session_state.num_quiz_questions,
            step=3,
            key="step6_initial_quiz_count_slider"
        )
        st.session_state.num_quiz_questions = num_q
        
        if st.button("Generate Quiz ✨" if lang == "English" else "توليد الاختبار ✨", type="primary", use_container_width=True):
            with st.spinner("Generating end-of-workshop comprehensive quiz..." if lang == "English" else "جاري توليد اختبار نهاية ورشة العمل التقييمي..."):
                try:
                    quiz = generate_quiz(
                        title=st.session_state.chosen_title,
                        outline=st.session_state.plan_result["outline"],
                        content=st.session_state.content_result,
                        min_questions=st.session_state.num_quiz_questions,
                        max_questions=st.session_state.num_quiz_questions
                    )
                    st.session_state.quiz_result = quiz
                    st.session_state.quiz_result_ar = None
                    st.session_state.quiz_approved = False
                    st.session_state.quiz_approved_snapshot = None
                    st.rerun()
                except Exception as e:
                    st.error(f"Error during quiz generation: {e}")

    st.markdown("---")
    
    col_nav1, col_nav2 = st.columns(2)
    with col_nav1:
        if st.button(T[lang]["btn_back"], use_container_width=True):
            st.session_state.step = 5
            st.rerun()
    with col_nav2:
        if st.button(T[lang]["btn_start_over"], type="secondary", use_container_width=True):
            init_state(force=True)
            clear_block_keys()
            st.rerun()



elif st.session_state.step == 7:
    st.header("Step 7: Final Quality Checklist 🔍" if lang == "English" else "الخطوة 7: قائمة الجودة النهائية 🔍")
    st.write(
        "Reviews the plan, content, labs, and the approved quiz TOGETHER, "
        "and gates the final export downloads on passing this check."
        if lang == "English" else
        "تراجع الخطة والمحتوى واللابات والاختبار المعتمد مع بعض، وتشترط اجتياز "
        "هذا الفحص قبل فتح ملفات التصدير النهائية."
    )

    # Safety net: this page assumes the quiz is already approved (that's
    # the only way to reach it, via the Step 6 button) — if session state
    # somehow ended up here without an approved quiz (e.g. a stale link
    # after starting over), bounce back to Step 6 instead of showing a
    # broken page.
    if not (st.session_state.quiz_approved and st.session_state.quiz_approved_snapshot):
        st.warning(
            "The quiz isn't approved yet — go back to Step 6 and approve it first."
            if lang == "English" else
            "الاختبار لسه ما اعتُمد — ارجعي للخطوة 6 واعتمديه أول."
        )
        if st.button("⬅ Back to Quiz" if lang == "English" else "⬅ رجوع للاختبار", key="step7_bounce_back"):
            st.session_state.step = 6
            st.rerun()
    else:
        # --- Step 7: Final Quality Checklist ---
        # Only runs once the quiz is approved — reviewing a still-changing
        # draft quiz would produce feedback on a version that might not
        # even ship. Reviews the plan, content, labs, and the FINAL
        # approved quiz together, not any one piece alone.
        if st.session_state.quiz_approved and st.session_state.quiz_approved_snapshot:
            st.markdown(
                "### 🔍 Step 7: Final Quality Checklist"
                if lang == "English" else
                "### 🔍 الخطوة 7: قائمة الجودة النهائية"
            )
            st.caption(
                "Reviews the plan, content, labs, and quiz TOGETHER before export."
                if lang == "English" else
                "تراجع الخطة والمحتوى واللابات والاختبار مع بعض قبل التصدير."
            )

            if st.button(
                "🔍 Run Quality Checklist" if lang == "English" else "🔍 تشغيل قائمة الجودة",
                key="run_quality_checklist_btn"
            ):
                with st.spinner(
                    "Reviewing the full workshop package..." if lang == "English"
                    else "جاري مراجعة الحزمة الكاملة..."
                ):
                    try:
                        result = run_quality_checklist(
                            title=st.session_state.chosen_title,
                            plan=st.session_state.plan_result,
                            content=st.session_state.content_result,
                            labs_result=st.session_state.labs_result,
                            quiz_result=st.session_state.quiz_approved_snapshot,
                        )
                        st.session_state.quality_result = result
                        # Pre-fill each failing category's feedback box with
                        # Step 7's own issues, as EDITABLE starting text.
                        # Nothing is sent anywhere from here — the trainer
                        # reviews/edits it below and must click Apply
                        # themselves before any revision happens.
                        feedback_texts = {}
                        for check in result.get("checks", []):
                            if check["status"] == "fail" and check["issues"]:
                                feedback_texts[check["category"]] = "\n".join(
                                    f"- {i}" for i in check["issues"]
                                )
                        st.session_state.quality_feedback_texts = feedback_texts
                        st.rerun()
                    except Exception as e:
                        st.error(f"Quality checklist failed: {e}")

            if st.session_state.quality_result:
                q_result = st.session_state.quality_result
                overall = q_result.get("overall_status")
                if overall == "pass":
                    st.success(
                        "✅ PASS — package is ready for export."
                        if lang == "English" else
                        "✅ ناجح — الحزمة جاهزة للتصدير."
                    )
                else:
                    st.error(
                        "❌ FAIL — review the issues below before exporting."
                        if lang == "English" else
                        "❌ فشل — راجعي الملاحظات بالأسفل قبل التصدير."
                    )

                if q_result.get("automated_issues"):
                    st.markdown(
                        "**Automated structural issues (facts, not opinions):**"
                        if lang == "English" else
                        "**مشاكل بنيوية آلية (حقائق، مو رأي):**"
                    )
                    for issue in q_result["automated_issues"]:
                        st.warning(issue)

                category_labels = {
                    "plan": "Plan" if lang == "English" else "الخطة",
                    "content": "Content" if lang == "English" else "المحتوى",
                    "labs": "Labs" if lang == "English" else "اللابات",
                    "quiz": "Quiz" if lang == "English" else "الاختبار",
                    "cross_consistency": "Cross-consistency" if lang == "English" else "الاتساق العام",
                }
                revisable_categories = {"plan", "content", "labs", "quiz"}

                for check in q_result.get("checks", []):
                    cat = check["category"]
                    status_icon = "✅" if check["status"] == "pass" else "❌"
                    with st.expander(
                        f"{status_icon} {category_labels.get(cat, cat)}",
                        expanded=(check["status"] == "fail")
                    ):
                        if check["status"] == "pass":
                            st.write("No issues found." if lang == "English" else "ما فيه ملاحظات.")
                            continue

                        for issue in check["issues"]:
                            st.write(f"- {issue}")

                        if cat not in revisable_categories:
                            # cross_consistency has no single result of its
                            # own to revise — it's a signal pointing at
                            # content/labs/quiz; fix it via those boxes.
                            st.caption(
                                "Apply this by editing the relevant Content/Labs/Quiz "
                                "feedback box above instead." if lang == "English" else
                                "طبّقيها عن طريق تعديل مربع المحتوى/اللابات/الاختبار المناسب بالأعلى."
                            )
                            continue

                        # Editable, pre-filled feedback box — Step 7 drafts
                        # the starting text, but the trainer reads it, can
                        # edit or clear it completely, and NOTHING is sent
                        # to the AI until they click Apply themselves.
                        editable_key = f"quality_feedback_{cat}"
                        st.text_area(
                            "Feedback to apply (edit freely before sending):"
                            if lang == "English" else
                            "الملاحظة اللي بتُرسل (عدّليها بحرية قبل الإرسال):",
                            value=st.session_state.quality_feedback_texts.get(cat, ""),
                            key=editable_key,
                        )

                        # Content is special: it can hold dozens of slides,
                        # and sending the WHOLE thing for one broad rewrite
                        # is unreliable — the model struggles to leave 18
                        # untouched slides alone while precisely fixing 2.
                        # Instead, the trainer picks exactly which slide(s)
                        # the feedback applies to, and only those get
                        # revised (via revise_single_slide), one at a time.
                        #
                        # Labs get the same treatment for the same reason,
                        # just at a smaller scale (usually 2-4 labs, not
                        # 20) — a fix aimed at one lab's missing setup
                        # instructions shouldn't risk the model quietly
                        # rewriting the other lab's notebook cells too.
                        selected_slide_indices = []
                        selected_lab_indices = []
                        if cat == "content":
                            all_slides = st.session_state.content_result.get("slides", [])
                            slide_options = {
                                i: f"#{s.get('slide_number', i)} — {s.get('slide_title', '(untitled)')} "
                                   f"[{s.get('section', '')}]"
                                for i, s in enumerate(all_slides)
                            }
                            # Auto-preselect slides whose title is quoted in
                            # Step 7's issue text (e.g. "the section on
                            # 'The Prompt...'") — a starting suggestion the
                            # trainer can freely adjust, not a final decision.
                            issue_text_combined = " ".join(check["issues"]).lower()
                            preselected = [
                                i for i, s in enumerate(all_slides)
                                if s.get("slide_title") and s["slide_title"].lower()[:15] in issue_text_combined
                            ]
                            selected_slide_indices = st.multiselect(
                                "Which slide(s) does this feedback apply to? (only these get revised)"
                                if lang == "English" else
                                "هالملاحظة تخص أي سلايد؟ (بس هذي بتتعدل)",
                                options=list(slide_options.keys()),
                                default=preselected,
                                format_func=lambda i: slide_options[i],
                                key=f"quality_slide_pick_{cat}",
                            )
                        elif cat == "labs":
                            all_labs = st.session_state.labs_result.get("labs", []) if st.session_state.labs_result else []
                            lab_options = {
                                i: f"[{lab.get('lab_type', '?')}] {lab.get('title', '(untitled)')} "
                                   f"(covers: {', '.join(lab.get('covers_sections', []))})"
                                for i, lab in enumerate(all_labs)
                            }
                            issue_text_combined = " ".join(check["issues"]).lower()
                            preselected_labs = [
                                i for i, lab in enumerate(all_labs)
                                if lab.get("title") and lab["title"].lower()[:15] in issue_text_combined
                            ]
                            selected_lab_indices = st.multiselect(
                                "Which lab(s) does this feedback apply to? (only these get revised)"
                                if lang == "English" else
                                "هالملاحظة تخص أي لاب؟ (بس هذا بيتعدل)",
                                options=list(lab_options.keys()),
                                default=preselected_labs,
                                format_func=lambda i: lab_options[i],
                                key=f"quality_lab_pick_{cat}",
                            )

                        if st.button(
                            f"Apply feedback & regenerate {category_labels.get(cat, cat)}"
                            if lang == "English" else
                            f"طبّقي الملاحظة وأعيدي توليد {category_labels.get(cat, cat)}",
                            key=f"apply_quality_fix_{cat}"
                        ):
                            feedback_text = st.session_state.get(editable_key, "").strip()
                            if not feedback_text:
                                st.warning(
                                    "Feedback box is empty — nothing to apply."
                                    if lang == "English" else
                                    "مربع الملاحظة فاضي — ما فيه شي نرسله."
                                )
                            elif cat == "content" and not selected_slide_indices:
                                st.warning(
                                    "Pick at least one slide above — nothing was selected to revise."
                                    if lang == "English" else
                                    "اختاري سلايد واحد على الأقل بالأعلى — ما فيه شي محدد نعدله."
                                )
                            elif cat == "labs" and not selected_lab_indices:
                                st.warning(
                                    "Pick at least one lab above — nothing was selected to revise."
                                    if lang == "English" else
                                    "اختاري لاب واحد على الأقل بالأعلى — ما فيه شي محدد نعدله."
                                )
                            else:
                                with st.spinner(
                                    "Applying feedback..." if lang == "English"
                                    else "جاري تطبيق الملاحظة..."
                                ):
                                    try:
                                        title_ctx = st.session_state.chosen_title
                                        if cat == "plan":
                                            st.session_state.plan_result = revise_output(
                                                st.session_state.plan_result, feedback_text,
                                                f"The workshop plan for '{title_ctx}'"
                                            )
                                            st.session_state.plan_result_ar = None
                                        elif cat == "content":
                                            # Revise ONLY the selected slides,
                                            # one at a time — not the whole
                                            # slides array — so the fix is
                                            # precise and the other slides
                                            # are guaranteed untouched.
                                            slides = st.session_state.content_result.get("slides", [])
                                            for idx in selected_slide_indices:
                                                slides[idx] = revise_single_slide(slides[idx], feedback_text)
                                            st.session_state.content_result["slides"] = slides
                                            st.session_state.content_result_ar = None
                                        elif cat == "labs":
                                            # Same principle as content:
                                            # revise ONLY the selected lab(s),
                                            # one at a time, leaving any other
                                            # lab's notebook cells untouched.
                                            labs = st.session_state.labs_result.get("labs", [])
                                            for idx in selected_lab_indices:
                                                labs[idx] = revise_output(
                                                    labs[idx], feedback_text,
                                                    f"One lab from the workshop '{title_ctx}' "
                                                    f"(titled '{labs[idx].get('title', '')}')"
                                                )
                                            st.session_state.labs_result["labs"] = labs
                                            st.session_state.labs_result_ar = None
                                        elif cat == "quiz":
                                            st.session_state.quiz_result = revise_output(
                                                st.session_state.quiz_result, feedback_text,
                                                f"Workshop quiz for '{title_ctx}'"
                                            )
                                            st.session_state.quiz_result_ar = None
                                            # A quiz edit invalidates the
                                            # frozen approved snapshot — the
                                            # trainer must re-approve before
                                            # exporting again.
                                            st.session_state.quiz_approved = False
                                            st.session_state.quiz_approved_snapshot = None

                                        # Any revision invalidates the last
                                        # quality check result — clear it so
                                        # the trainer re-runs Step 7 on the
                                        # new version rather than trusting a
                                        # stale verdict.
                                        st.session_state.quality_result = None
                                        st.session_state.quality_feedback_texts = {}
                                        st.success(
                                            "Applied. Re-run the checklist to verify."
                                            if lang == "English" else
                                            "تم التطبيق. أعيدي تشغيل القائمة للتأكد."
                                        )
                                        st.rerun()
                                    except Exception as e:
                                        st.error(f"Failed to apply feedback: {e}")

            st.markdown("---")

        # --- Downloads are gated on Step 7, but the gate has two tiers:
        # automated (structural) issues are FACTS — a section with no
        # Exports are ALWAYS available once the quiz is approved — Step 7
        # never blocks them. If a check has been run, its result shows as
        # an informational banner (warning if issues exist, success if
        # clean) so the trainer sees the state, but the decision to export
        # anyway is always theirs without needing to check a box.
        has_result = st.session_state.quality_result is not None
        if has_result:
            automated_issues_exist = bool(st.session_state.quality_result.get("automated_issues"))
            overall_status = st.session_state.quality_result.get("overall_status")
            if overall_status == "pass":
                st.success(
                    "✅ Step 7 passed — no issues found."
                    if lang == "English" else
                    "✅ اجتازت الخطوة 7 — ما فيه مشاكل."
                )
            elif automated_issues_exist:
                st.warning(
                    "⚠️ Step 7 found structural issues above (missing content, "
                    "broken quiz answers, etc.) — review them before exporting, "
                    "though nothing here stops you."
                    if lang == "English" else
                    "⚠️ الخطوة 7 لقت مشاكل بنيوية بالأعلى (محتوى ناقص، إجابات كويز "
                    "غلط، إلخ) — راجعيها قبل التصدير، بس ما فيه شي يوقفك."
                )
            else:
                st.warning(
                    "⚠️ Step 7 found judgment-based concerns above (quality/"
                    "consistency opinions, not structural errors) — review them "
                    "before exporting, though nothing here stops you."
                    if lang == "English" else
                    "⚠️ الخطوة 7 لقت ملاحظات حكمية بالأعلى (رأي عن الجودة/الاتساق، "
                    "مو أخطاء بنيوية) — راجعيها قبل التصدير، بس ما فيه شي يوقفك."
                )
        else:
            st.info(
                "You haven't run the Step 7 quality checklist above yet — "
                "exports work either way, but running it first is recommended."
                if lang == "English" else
                "لسه ما شغّلتي قائمة الجودة (الخطوة 7) بالأعلى — التصدير شغّال "
                "برضو، بس يُفضّل تشغيلها أول."
            )

        if st.session_state.quiz_approved and st.session_state.quiz_approved_snapshot:
            approved_quiz_dict = st.session_state.quiz_approved_snapshot

            # JSON Exporter — quiz portion comes from the approved snapshot
            export_payload = {
                "title": st.session_state.chosen_title,
                "audience": st.session_state.audience,
                "age": st.session_state.age,
                "duration": st.session_state.duration,
                "learning_objectives": st.session_state.plan_result.get("learning_objectives", []),
                "outline": st.session_state.plan_result.get("outline", []),
                "slides_content": st.session_state.content_result.get("slides", []),
                "labs": st.session_state.labs_result.get("labs", []) if st.session_state.labs_result else [],
                "quiz": approved_quiz_dict.get("quiz", {})
            }
            json_str = json.dumps(export_payload, indent=2, ensure_ascii=False)

            # Word doc Exporter — from the approved snapshot only
            docx_data = None
            try:
                save_quiz_docx(approved_quiz_dict, "quiz.docx")
                if os.path.exists("quiz.docx"):
                    with open("quiz.docx", "rb") as f:
                        docx_data = f.read()
            except Exception as e:
                st.warning(f"Could not build Word document (.docx): {e}")

            # Kahoot Exporter — from the approved snapshot only. Requires
            # Kahoot's own template file (KahootQuizTemplate.xlsx) to be
            # present in Backend/ — we fill it in, we don't build one from
            # scratch, since Kahoot's importer requires their exact template.
            kahoot_data = None
            kahoot_template_path = backend_dir / "KahootQuizTemplate.xlsx"
            try:
                if kahoot_template_path.exists():
                    export_quiz_to_kahoot_xlsx(
                        approved_quiz_dict,
                        template_path=str(kahoot_template_path),
                        output_path="kahoot_import.xlsx"
                    )
                    if os.path.exists("kahoot_import.xlsx"):
                        with open("kahoot_import.xlsx", "rb") as f:
                            kahoot_data = f.read()
                else:
                    st.warning(
                        f"Kahoot template not found at {kahoot_template_path}. Place "
                        "KahootQuizTemplate.xlsx in the Backend folder to enable this export."
                        if lang == "English" else
                        f"ملف قالب كاهوت غير موجود بمسار {kahoot_template_path}. حطي "
                        "KahootQuizTemplate.xlsx بمجلد Backend عشان يشتغل هذا التصدير."
                    )
            except Exception as e:
                st.warning(f"Could not build Kahoot import file: {e}")

            col_dl1, col_dl2, col_dl3 = st.columns(3)
            with col_dl1:
                st.download_button(
                    label=T[lang]["export_btn"],
                    data=json_str,
                    file_name="ai_workshop_studio_package.json",
                    mime="application/json",
                    use_container_width=True
                )
            with col_dl2:
                if docx_data:
                    st.download_button(
                        label="Download Quiz Word Doc (.docx) 📥" if lang == "English" else "تحميل أسئلة الاختبار كملف Word (.docx) 📥",
                        data=docx_data,
                        file_name="ai_workshop_quiz.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True
                    )
            with col_dl3:
                if kahoot_data:
                    st.download_button(
                        label="Download Kahoot Import (.xlsx) 🎮" if lang == "English" else "تحميل ملف كاهوت (.xlsx) 🎮",
                        data=kahoot_data,
                        file_name="kahoot_import.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True
                    )

            st.markdown("---")

            if st.button(
                "💾 Save to History" if lang == "English" else "💾 حفظ بالتاريخ",
                use_container_width=True,
                key="save_to_history_btn",
            ):
                try:
                    from workshop_db import save_workshop
                    new_id = save_workshop(
                        title=st.session_state.chosen_title,
                        audience=st.session_state.audience,
                        age=st.session_state.age,
                        duration=st.session_state.duration,
                        plan=st.session_state.plan_result,
                        content=st.session_state.content_result,
                        labs=st.session_state.labs_result,
                        quiz=approved_quiz_dict,
                    )
                    st.success(f"Saved! (id #{new_id})" if lang == "English" else f"تم الحفظ! (رقم #{new_id})")
                except Exception as e:
                    st.error(f"Could not save to history: {e}")

            st.markdown("---")

            # --- One-click "download everything" ZIP. Built fresh, in
            # memory, from the CURRENT approved data (quiz snapshot, live
            # content_result, live labs_result) — never from files that
            # might be sitting stale on disk from an earlier generation,
            # before any Step 7 revisions. ---
            if st.button(
                "📦 Prepare Full Package (.zip)" if lang == "English" else
                "📦 جهّزي الحزمة الكاملة (.zip)",
                use_container_width=True,
                key="build_full_zip_btn"
            ):
                with st.spinner(
                    "Building the full package..." if lang == "English" else
                    "جاري تجهيز الحزمة الكاملة..."
                ):
                    try:
                        zip_buffer = io.BytesIO()
                        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
                            # 1) The full JSON package (plan + content + labs + quiz)
                            zf.writestr("ai_workshop_studio_package.json", json_str)

                            # 2) Quiz — Word doc
                            if docx_data:
                                zf.writestr("Quiz/ai_workshop_quiz.docx", docx_data)

                            # 3) Quiz — Kahoot import (only if the template was found)
                            if kahoot_data:
                                zf.writestr("Quiz/kahoot_import.xlsx", kahoot_data)

                            # 4) Slides — PPTX, built fresh from the live content
                            try:
                                from pptx_export import export_to_pptx
                                import os
                                
                                pptx_opts = [
                                    "Generate default PowerPoint (Clean, Bold, etc.)" if lang == "English" else "توليد ملف بوربوينت افتراضي",
                                    "Upload a custom template to mimic its style and colors" if lang == "English" else "رفع قالب بوربوينت مخصص لمحاكاة النمط والألوان"
                                ]
                                
                                template_file_path = None
                                if st.session_state.get("pptx_generation_mode") == pptx_opts[1] and st.session_state.get("pptx_custom_template_file") is not None:
                                    uploaded = st.session_state.pptx_custom_template_file
                                    temp_path = "temp_zip_template.pptx"
                                    with open(temp_path, "wb") as tf_file:
                                        tf_file.write(uploaded.getbuffer())
                                    template_file_path = temp_path# Build a copy of slides with chosen styles
                                    import copy
                                    slides_to_export = copy.deepcopy(st.session_state.content_result.get("slides", []))
                                    if st.session_state.get("pptx_generation_mode") == pptx_opts[0]:
                                        chosen_style = st.session_state.get("slide_style", "Clean & Minimal")
                                        style_app = st.session_state.get("style_application", "All slides")
                                        
                                        for idx, slide_data in enumerate(slides_to_export):
                                            if style_app in ["All slides", "جميع الشرائح"]:
                                                slide_data["slide_style"] = chosen_style
                                            else:
                                                section_name = slide_data.get("section", "")
                                                outline = st.session_state.plan_result.get("outline", [])
                                                section_idx = 0
                                                for o_idx, o_sec in enumerate(outline):
                                                    if o_sec.get("section", "") == section_name:
                                                        section_idx = o_idx
                                                        break
                                                sec_style = st.session_state.get(f"style_section_{section_idx}", chosen_style)
                                                slide_data["slide_style"] = sec_style
                                    
                                    pptx_tmp_path = export_to_pptx(
                                        title=st.session_state.chosen_title,
                                        slides=slides_to_export,
                                        template_path=template_file_path,
                                        include_images=st.session_state.get("pptx_include_images", True)
                                    )
                                
                                if template_file_path and os.path.exists(template_file_path):
                                    try:
                                        os.remove(template_file_path)
                                    except Exception:
                                        pass
                                with open(pptx_tmp_path, "rb") as f:
                                    zf.writestr(
                                        f"Slides/{st.session_state.chosen_title.replace(' ', '_')}.pptx",
                                        f.read(),
                                    )
                            except Exception as e:
                                st.warning(f"Could not include slides (.pptx): {e}")

                            # 5) Labs — real .ipynb notebooks for "coding" labs
                            # (built fresh in memory, not read from generated_labs/
                            # on disk, so a Step 7 revision is always reflected),
                            # and a simple Q&A markdown file for "conceptual" labs.
                            if st.session_state.labs_result:
                                from notebook_builder import build_notebook, sanitize_filename
                                import nbformat
                                for lab in st.session_state.labs_result.get("labs", []):
                                    base_name = sanitize_filename(lab.get("title", "lab"))
                                    if lab.get("lab_type") == "coding":
                                        try:
                                            trainee_nb = build_notebook(lab.get("trainee_notebook_cells", []))
                                            solution_nb = build_notebook(lab.get("solution_notebook_cells", []))
                                            zf.writestr(
                                                f"Labs/{base_name}_trainee.ipynb",
                                                nbformat.writes(trainee_nb),
                                            )
                                            zf.writestr(
                                                f"Labs/{base_name}_solution.ipynb",
                                                nbformat.writes(solution_nb),
                                            )
                                        except Exception as e:
                                            st.warning(f"Could not include lab '{lab.get('title')}': {e}")
                                    elif lab.get("lab_type") == "conceptual":
                                        qa_lines = [f"# {lab.get('title', 'Lab')}", "", lab.get("instructions", ""), ""]
                                        for i, q in enumerate(lab.get("questions", []), start=1):
                                            qa_lines.append(f"**Q{i}.** {q.get('question', '')}")
                                            qa_lines.append(f"*Answer:* {q.get('answer', '')}")
                                            qa_lines.append("")
                                        zf.writestr(f"Labs/{base_name}.md", "\n".join(qa_lines))

                        st.session_state["full_zip_bytes"] = zip_buffer.getvalue()
                        st.success(
                            "Package ready — click below to download."
                            if lang == "English" else
                            "الحزمة جاهزة — اضغطي التحميل بالأسفل."
                        )
                    except Exception as e:
                        st.error(f"Failed to build the full package: {e}")

            if st.session_state.get("full_zip_bytes"):
                st.download_button(
                    label="⬇️ Download Full Package (.zip)" if lang == "English" else
                          "⬇️ تحميل الحزمة الكاملة (.zip)",
                    data=st.session_state["full_zip_bytes"],
                    file_name=f"{st.session_state.chosen_title.replace(' ', '_')}_package.zip",
                    mime="application/zip",
                    type="primary",
                    use_container_width=True,
                )

    st.markdown("---")

    col_nav1, col_nav2 = st.columns(2)
    with col_nav1:
        if st.button(
            "⬅ Back to Quiz" if lang == "English" else "⬅ رجوع للاختبار",
            use_container_width=True,
            key="step7_back_btn"
        ):
            st.session_state.step = 6
            st.rerun()
    with col_nav2:
        if st.button(
            T[lang]["btn_start_over"],
            type="secondary",
            use_container_width=True,
            key="step7_start_over_btn"
        ):
            init_state(force=True)
            clear_block_keys()
            st.rerun()


# Save state to disk at the end of every run
save_state_to_disk()